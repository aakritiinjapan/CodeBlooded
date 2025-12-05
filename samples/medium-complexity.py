// """
// Enhanced PowerPoint Content Processor with Simple Text Extraction
// Implements the new workflow:
// 1. Remove useless pages (thank you, appendix, etc)
// 2. Separate slides into simple, complex and table slides
// 3. Extract text from simple slides using simple text extraction
// 4. Extract image descriptions using Gemini Vision for complex slides
// 5. Process tables with Gemini for CSV extraction
// """

import os
import re
import base64
import subprocess
import tempfile
import traceback
import shutil
from pathlib import Path
from typing import Dict, List, Any, Optional, Union, Tuple
from datetime import datetime
import pandas as pd

# PowerPoint processing
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

# Cross-platform utilities
from platform_utils import (
    get_platform_capabilities, 
    is_windows, 
    has_windows_com, 
    has_libreoffice
)
from office_utils import (
    ppt_to_pdf, 
    OfficeConversionError,
    log_conversion_capabilities
)
from powerpoint_utils import (
    PowerPointSlideExporter,
    export_powerpoint_slide,
    create_powerpoint_pdf
)

# Windows-specific imports (conditional)
try:
    import comtypes.client
    import pythoncom
    WINDOWS_COM_AVAILABLE = True
except ImportError:
    WINDOWS_COM_AVAILABLE = False
    print("Warning: Windows COM libraries not available. PowerPoint processing will use basic text extraction only.")

from file_loader import TimestampJSONEncoder, save_intermediate_file
from utils import call_gemini, encode_file

logger = logging.getLogger(__name__)

class PowerPointProcessor:
    """Enhanced PowerPoint processor with simple text extraction workflow"""
    
    def __init__(self):
        self.file_path = ""
        self.full_file_path = None
        self.processed_slides = {}
        self.table_slides = {}
        self.regular_slides = {}
        
        # Create output directories
        self.table_slides_dir = Path("table_slides")
        self.table_slides_dir.mkdir(exist_ok=True)
        
        # Create CSV output directory for 3-agent pipeline processing
        self.csv_output_dir = Path("csv_for_pipeline")
        self.csv_output_dir.mkdir(exist_ok=True)
        
        # Initialize slide exporter
        self.slide_exporter = PowerPointSlideExporter()
        
        # Log platform capabilities for debugging
        log_conversion_capabilities()
    
    def __del__(self):
        """Ensure cleanup on object destruction"""
        if hasattr(self, 'slide_exporter'):
            self.slide_exporter.cleanup_temp_files()
        
    def process_powerpoint_file(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Main processing function implementing the new simple OCR workflow
        
        Args:
            file_path: Path to the PowerPoint file
            
        Returns:
            Dictionary with processed content ready for 3-agent pipeline
        """
        try:
            if isinstance(file_path, str):
                file_path = Path(file_path)
                
            self.file_path = file_path.name
            self.full_file_path = file_path  # Store full path for slide export
            print(f">> [STEP 0] Starting PowerPoint processing: {file_path.name}")
            print(f">> File exists: {file_path.exists()}")
            print(f">> File is file: {file_path.is_file()}")
            print(f">> File absolute path: {file_path.absolute()}")
            
            # Step 1: Analyze slides and remove useless pages
            print(f">> [STEP 1] Analyzing slides and filtering useless pages...")
            valid_slides, removed_slides = self._filter_useless_slides(file_path)
            print(f">> [STEP 1] Filtered {len(removed_slides)} useless slides, processing {len(valid_slides)} slides")
            print(f">> [STEP 1] Valid slides: {valid_slides}")
            print(f">> [STEP 1] Removed slides: {removed_slides}")
            
            # Step 2: Separate slides by complexity and content type
            print(f">> [STEP 2] Separating slides by complexity and content type...")
            table_slides, complex_slides, simple_slides = self._separate_slides_by_complexity(file_path, valid_slides)
            print(f">> [STEP 2] Found {len(table_slides)} table slides, {len(complex_slides)} complex slides, and {len(simple_slides)} simple slides")
            print(f">> [STEP 2] Table slides: {table_slides}")
            print(f">> [STEP 2] Complex slides: {complex_slides}")
            print(f">> [STEP 2] Simple slides: {simple_slides}")
            
            # Step 3: Process table slides separately with Gemini
            print(f">> [STEP 3] Processing table slides with Gemini...")
            processed_tables = self._process_table_slides(file_path, table_slides)
            print(f">> [STEP 3] Processed {len(processed_tables)} table slides")
            
            # Step 4: Process complex slides separately with Gemini Vision
            print(f">> [STEP 4] Processing complex slides separately with Gemini Vision...")
            complex_slide_descriptions = self._process_complex_slides_with_gemini(file_path, complex_slides)
            print(f">> [STEP 4] Processed {len(complex_slide_descriptions)} complex slides")
            
            # Step 5: Extract text from simple slides using simple text extraction
            print(f">> [STEP 5] Extracting text from simple slides using simple text extraction...")
            markdown_content = self._process_simple_slides_with_text_extraction(file_path, simple_slides)
            print(f">> [STEP 5] Generated markdown content: {len(markdown_content) if markdown_content else 0} characters")
            
            # Save the raw text extraction markdown
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            raw_text_filename = f"{file_path.stem}_raw_text_extraction_{timestamp}.md"
            raw_text_path = Path(raw_text_filename)

            if markdown_content:
                print(f">> [STEP 5.1] Saving raw text extraction markdown to: {raw_text_filename}")
                with open(raw_text_path, 'w', encoding='utf-8') as f:
                    f.write(f"# Raw Text Extraction from {file_path.name}\n\n")
                    f.write(f"*Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*\n\n")
                    f.write(f"## Simple Slides - Text Extraction\n\n")
                    f.write(markdown_content)
                    f.write(f"\n\n## Complex Slides (Image-based extraction)\n\n")
                    if complex_slide_descriptions and len(complex_slide_descriptions) > 0:
                        for slide_num, description in complex_slide_descriptions.items():
                            f.write(f"### {file_path.name}:slide-{slide_num}\n")
                            f.write(description.strip() + "\n\n")
                    else:
                        f.write("*Complex slides identified but not yet processed - placeholder for future image extraction*\n")
                print(f">> [STEP 5.1] Raw text extraction markdown saved successfully")
            else:
                print(f">> [STEP 5.1] No markdown content to save")
                raw_text_path = None
            
            # Step 6: Process CSVs through 3-agent pipeline and combine results
            print(f">> [STEP 6] Processing CSVs through 3-agent pipeline and combining results...")
            final_result = self._process_csvs_and_combine_results(
                markdown_content, 
                processed_tables, 
                file_path.name, 
                raw_text_path, 
                None,  # No PDF file path needed for simple text extraction
                complex_slide_descriptions, 
                simple_slides
            )
            print(f">> [STEP 6] Final result type: {type(final_result)}")
            print(f">> [STEP 6] Final result keys: {list(final_result.keys()) if final_result else 'None'}")
            
            print(f">> [COMPLETED] PowerPoint processing for {file_path.name}")
            return final_result
            
        except Exception as e:
            print(f">> [ERROR] Exception in process_powerpoint_file: {e}")
            print(f">> [ERROR] Exception type: {type(e).__name__}")
            import traceback
            print(f">> [ERROR] Full traceback:")
            traceback.print_exc()
            return {
                'success': False,
                'error': str(e),
                'exception_type': type(e).__name__
            }
    
    def _filter_useless_slides(self, file_path: Path) -> Tuple[List[int], List[int]]:
        """Filter out useless slides like thank you, appendix, etc. Keep title slides with important info."""
        try:
            prs = Presentation(str(file_path))
            valid_slides = []
            removed_slides = []
            
            # More specific patterns - avoid removing title slides with dates/project info
            skip_patterns = [
                'thank you',
                'thanks',
                'appendix',
                'questions',
                'q&a',
                'contact us',
                'contact information',
                'contact:'
            ]
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_text = ""
                
                # Extract all text from slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += " " + shape.text.strip().lower()
                
                # Check if slide should be skipped - be more aggressive with thank you/appendix
                should_skip = False
                
                # Skip if slide contains these patterns, regardless of other content
                for pattern in skip_patterns:
                    if pattern in slide_text:
                        should_skip = True
                        logger.debug(f"Skipping slide {slide_idx} - matches pattern '{pattern}'")
                        break
                
                if should_skip:
                    removed_slides.append(slide_idx)
                else:
                    valid_slides.append(slide_idx)
            
            return valid_slides, removed_slides
            
        except Exception as e:
            logger.error(f"Error filtering slides: {e}")
            # Fallback: process all slides
            prs = Presentation(str(file_path))
            return list(range(1, len(prs.slides) + 1)), []
    
    def _separate_table_slides(self, file_path: Path, valid_slides: List[int]) -> Tuple[List[int], List[int]]:
        """Separate slides that contain tables from regular content slides"""
        try:
            prs = Presentation(str(file_path))
            table_slides = []
            regular_slides = []
            
            for slide_idx in valid_slides:
                slide = prs.slides[slide_idx - 1]  # Convert to 0-based index
                has_table = False
                
                # Check for table shapes
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        has_table = True
                        logger.debug(f"Found table in slide {slide_idx}")
                        break
                
                if has_table:
                    table_slides.append(slide_idx)
                else:
                    regular_slides.append(slide_idx)
            
            return table_slides, regular_slides
            
        except Exception as e:
            logger.error(f"Error separating table slides: {e}")
            # Fallback: treat all as regular slides
            return [], valid_slides

    def _is_complex_slide(self, slide) -> bool:
        """
        Determine if a slide is complex and should use image-based extraction.
        
        Complex slides have:
        - Multiple visual graphics (charts, diagrams, images)
        - Layers, overlays, or connectors  
        - Diagrams combined with tables
        - High content object count (excluding placeholders)
        
        Returns True if slide should use image-based extraction
        """
        try:
            # Separate content objects from placeholder/infrastructure objects
            content_objects = []
            placeholder_objects = []
            
            for shape in slide.shapes:
                shape_name = getattr(shape, 'name', '').lower()
                
                # Identify placeholder shapes (footer, slide number, etc.)
                is_placeholder = (
                    'footer' in shape_name or
                    'slide number' in shape_name or
                    'placeholder' in shape_name
                )
                
                # Also check if it's actually a placeholder shape type
                try:
                    if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                        is_placeholder = True
                except:
                    pass  # Not a placeholder shape
                
                if is_placeholder:
                    placeholder_objects.append(shape)
                else:
                    content_objects.append(shape)
            
            # Analyze only content objects (excluding placeholders)
            content_obj_types = [shape.shape_type for shape in content_objects]
            num_content_objects = len(content_obj_types)
            
            # Count specific object types in content
            num_tables = content_obj_types.count(19)  # Table shape type
            num_text_boxes = content_obj_types.count(17)  # Text box shape type  
            num_auto_shapes = content_obj_types.count(1)  # Auto shape type
            num_pictures = content_obj_types.count(13)  # Picture shape type
            num_lines = content_obj_types.count(9)  # Line/connector shape type
            
            # Visual graphics: pictures, complex auto shapes (excluding simple text containers)
            visual_graphics = []
            for shape in content_objects:
                shape_type = shape.shape_type
                shape_name = getattr(shape, 'name', '').lower()
                
                # Pictures are always visual
                if shape_type == 13:  # Picture
                    visual_graphics.append(shape)
                # Lines/connectors are structural graphics
                elif shape_type == 9:  # Line
                    visual_graphics.append(shape)
                # Auto shapes that are not just text containers
                elif shape_type == 1:  # Auto shape
                    # Check if it's a simple text container (rectangle with just text)
                    has_text = hasattr(shape, 'text') and shape.text.strip()
                    is_simple_rect = 'rectangle' in shape_name or 'rect' in shape_name
                    
                    # If it's not a simple text rectangle, consider it visual
                    if not (is_simple_rect and has_text and len(shape.text.strip()) > 10):
                        visual_graphics.append(shape)
            
            # Check for connectors
            has_connectors = any("connector" in getattr(shape, 'name', '').lower() for shape in content_objects)
            
            # Check for grouped objects (complex layouts)
            has_groups = any(hasattr(shape, 'shapes') for shape in content_objects)
            
            # Complexity criteria (now more specific)
            is_complex = (
                num_content_objects > 6 or  # High content object count (excluding placeholders)
                len(visual_graphics) > 2 or  # Multiple visual graphics
                has_connectors or  # Has connector lines
                has_groups or  # Has grouped objects
                num_pictures > 0 or  # Any pictures/images
                (num_tables > 0 and len(visual_graphics) > 0)  # Mixed tables + visual elements
            )
            
            if is_complex:
                logger.info(f"Slide classified as COMPLEX: {num_content_objects} content objects, "
                          f"{len(visual_graphics)} visual graphics, "
                          f"connectors={has_connectors}, groups={has_groups}")
            else:
                logger.debug(f"Slide classified as SIMPLE: {num_content_objects} content objects, "
                           f"{len(visual_graphics)} visual graphics")
            
            return is_complex
            
        except Exception as e:
            logger.error(f"Error analyzing slide complexity: {e}")
            # Default to simple if analysis fails
            return False

    def _separate_slides_by_complexity(self, file_path: Path, valid_slides: List[int]) -> Tuple[List[int], List[int], List[int]]:
        """
        Separate slides into table slides, complex slides (for image extraction), and simple slides (for text extraction).
        
        Returns:
            table_slides: Slides with isolated tables (process with Gemini table extraction)
            complex_slides: Slides with diagrams/complex layouts (process with image extraction)  
            simple_slides: Plain text/simple slides (process with text extraction)
        """
        try:
            prs = Presentation(str(file_path))
            table_slides = []
            complex_slides = []
            simple_slides = []

            for slide_idx in valid_slides:
                slide = prs.slides[slide_idx - 1]  # Convert to 0-based index

                # Check if slide is complex first
                if self._is_complex_slide(slide):
                    complex_slides.append(slide_idx)
                    logger.info(f"Slide {slide_idx} → COMPLEX (image-based extraction)")
                    continue

                # Check for both table and text box
                has_table = False
                has_textbox = False
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        has_table = True
                    if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                        # Check for non-empty text
                        if hasattr(shape, 'text') and shape.text.strip():
                            has_textbox = True

                # If both table and text box, classify as complex
                if has_table and has_textbox:
                    complex_slides.append(slide_idx)
                    logger.info(f"Slide {slide_idx} → COMPLEX (table + text box)")
                elif has_table:
                    table_slides.append(slide_idx)
                    logger.info(f"Slide {slide_idx} → TABLE (Gemini table extraction)")
                else:
                    simple_slides.append(slide_idx)
                    logger.info(f"Slide {slide_idx} → SIMPLE (text extraction)")

            logger.info(f"Slide classification complete: {len(table_slides)} table, "
                       f"{len(complex_slides)} complex, {len(simple_slides)} simple")

            return table_slides, complex_slides, simple_slides

        except Exception as e:
            logger.error(f"Error separating slides by complexity: {e}")
            # Fallback: use original table separation logic
            table_slides, regular_slides = self._separate_table_slides(file_path, valid_slides)
            return table_slides, [], regular_slides
    
    def _process_table_slides(self, file_path: Path, table_slides: List[int]) -> Dict[str, Any]:
        """Process slides containing tables using Gemini only (simplified approach)"""
        processed_tables = {}
        
        if not table_slides:
            return processed_tables
        
        try:
            print(f">> [TABLE] Processing {len(table_slides)} table slides with Gemini only...")
            prs = Presentation(str(file_path))
            
            for slide_idx in table_slides:
                print(f">> [TABLE] Processing table slide {slide_idx} with Gemini...")
                slide = prs.slides[slide_idx - 1]
                
                # Extract all text from the slide (including table text)
                slide_text = self._extract_text_from_slide(slide)
                print(f">> [TABLE] Extracted {len(slide_text)} characters of text from table slide {slide_idx}")
                
                # Process the table slide with Gemini (same as complex slides)
                description = self._convert_slide_to_image_and_process(slide, slide_idx, file_path.name)
                print(f">> [TABLE] Generated description for table slide {slide_idx}: {len(description) if description else 0} characters")
                
                # Store the processed content
                processed_tables[slide_idx] = {
                    'slide_number': slide_idx,
                    'description': description,
                    'extracted_text': slide_text
                }
                
                print(f">> [TABLE] Table slide {slide_idx} processing completed")
                
            print(f">> [TABLE] All {len(table_slides)} table slides processed with Gemini")
            return processed_tables
            
        except Exception as e:
            print(f">> [TABLE] Error processing table slides: {e}")
            logger.error(f"Error processing table slides: {e}")
            return {}

    def _process_complex_slides_with_gemini(self, file_path: Path, complex_slides: List[int]) -> Dict[str, str]:
        """
        Process complex slides separately with Gemini Vision API.
        
        Returns a dictionary mapping slide numbers to their Gemini descriptions
        with proper metadata formatting.
        """
        if not complex_slides:
            print(f">> [COMPLEX] No complex slides to process")
            return {}
            
        try:
            print(f">> [COMPLEX] Processing {len(complex_slides)} complex slides separately with Gemini Vision...")
            print(f">> [COMPLEX] Complex slides list: {complex_slides}")
            
            complex_slide_descriptions = {}
            prs = Presentation(str(file_path))
            
            for slide_idx in complex_slides:
                print(f">> [COMPLEX] Starting processing of complex slide {slide_idx}")
                slide = prs.slides[slide_idx - 1]
                
                # Analyze slide complexity details for logging
                obj_types = [shape.shape_type for shape in slide.shapes]
                num_objects = len(obj_types)
                
                # Get shape type counts for detailed analysis
                shape_counts = {}
                for shape_type in set(obj_types):
                    shape_counts[shape_type] = obj_types.count(shape_type)
                
                print(f">> [COMPLEX] Slide {slide_idx} analysis:")
                print(f"   - Total objects: {num_objects}")
                print(f"   - Shape types: {list(shape_counts.keys())}")
                print(f"   - Shape counts: {shape_counts}")
                
                # Convert slide to image and process with Gemini
                print(f">> [COMPLEX] Converting slide {slide_idx} to image and processing with Gemini...")
                slide_description = self._convert_slide_to_image_and_process(slide, slide_idx, file_path.name)
                print(f">> [COMPLEX] Slide {slide_idx} description generated: {len(slide_description) if slide_description else 0} characters")
                
                # Format with metadata
                metadata = f"{file_path.name}:slide-{slide_idx}"
                formatted_description = f"## {metadata}\n\n{slide_description}"
                
                complex_slide_descriptions[slide_idx] = formatted_description
                print(f">> [COMPLEX] Slide {slide_idx} processing completed successfully")
            
            print(f">> [COMPLEX] All {len(complex_slides)} complex slides processed")
            return complex_slide_descriptions
            
        except Exception as e:
            print(f">> [COMPLEX] Error processing complex slides: {e}")
            import traceback
            print(f">> [COMPLEX] Full traceback:")
            traceback.print_exc()
            return {}
            logger.error(f"Error processing complex slides: {e}")
            return {}

    def _process_simple_slides_with_text_extraction(self, file_path: Path, simple_slides: List[int]) -> str:
        """
        Process ONLY simple slides with basic text extraction, preserving source attribution.
        This replaces the previous OCR approach with a simple text extraction method.
        """
        if not simple_slides:
            print(f">> [TEXT] No simple slides to process")
            return ""
        
        try:
            print(f">> [TEXT] Processing {len(simple_slides)} simple slides with text extraction")
            print(f">> [TEXT] Simple slides list: {simple_slides}")
            
            # Use the existing fallback text extraction method
            extracted_content = self._extract_fallback_text_from_slides(str(file_path), simple_slides)
            
            if extracted_content:
                print(f">> [TEXT] Successfully extracted text from {len(simple_slides)} slides")
                print(f">> [TEXT] Total text length: {len(extracted_content)} characters")
                return extracted_content
            else:
                print(f">> [TEXT] No text content extracted from simple slides")
                return ""
                
        except Exception as e:
            logger.error(f"Error processing simple slides with text extraction: {e}")
            print(f">> [TEXT] Error in text extraction: {e}")
            return ""
    
    def _create_single_slide_presentation(self, file_path: Path, slide_number: int) -> str:
        """Create a PPTX with only one specified slide using PowerPoint COM automation"""
        if not WINDOWS_COM_AVAILABLE:
            print(f">> [WARNING] Windows COM not available. Skipping single slide creation for slide {slide_number}")
            return ""
            
        powerpoint = None
        original_presentation = None
        new_presentation = None
        
        try:
            # Initialize COM
            pythoncom.CoInitialize()
            
            # Create PowerPoint application
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = 1
            powerpoint.WindowState = 2  # Minimize window
            
            # Open original presentation
            abs_path = os.path.abspath(str(file_path))
            original_presentation = powerpoint.Presentations.Open(abs_path, ReadOnly=True)
            
            # Create new presentation
            new_presentation = powerpoint.Presentations.Add()
            
            # Remove default slide if present
            if new_presentation.Slides.Count > 0:
                new_presentation.Slides[1].Delete()
            
            # Copy the specific slide (PowerPoint uses 1-based indexing)
            slide_to_copy = original_presentation.Slides[slide_number]
            slide_to_copy.Copy()
            new_presentation.Slides.Paste()
            
            # Create timestamp for unique filename
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            single_slide_path = f"temp_slide_{slide_number}_{timestamp}.pptx"
            
            # Save the new presentation
            new_presentation.SaveAs(os.path.abspath(single_slide_path))
            
            logger.info(f"✅ Created single-slide presentation: {single_slide_path}")
            return single_slide_path
            
        except Exception as e:
            logger.error(f"Error creating single slide presentation: {e}")
            return ""
        
        finally:
            # Ensure proper cleanup
            try:
                if new_presentation:
                    new_presentation.Close()
            except:
                pass
            try:
                if original_presentation:
                    original_presentation.Close()
            except:
                pass
            try:
                if powerpoint:
                    powerpoint.Quit()
            except:
                pass
            try:
                pythoncom.CoUninitialize()
            except:
                pass
    
    def _process_complex_slides_with_gemini(self, file_path: Path, complex_slides: List[int]) -> Dict[str, str]:
        """
        Process complex slides separately with Gemini Vision API.
        
        Returns a dictionary mapping slide numbers to their Gemini descriptions
        with proper metadata formatting.
        """
        if not complex_slides:
            print(f">> [COMPLEX] No complex slides to process")
            return {}
            
        try:
            print(f">> [COMPLEX] Processing {len(complex_slides)} complex slides separately with Gemini Vision...")
            print(f">> [COMPLEX] Complex slides list: {complex_slides}")
            
            complex_slide_descriptions = {}
            prs = Presentation(str(file_path))
            
            for slide_idx in complex_slides:
                print(f">> [COMPLEX] Starting processing of complex slide {slide_idx}")
                slide = prs.slides[slide_idx - 1]
                
                # Analyze slide complexity details for logging
                obj_types = [shape.shape_type for shape in slide.shapes]
                num_objects = len(obj_types)
                
                # Get shape type counts for detailed analysis
                shape_counts = {}
                for shape_type in set(obj_types):
                    shape_counts[shape_type] = obj_types.count(shape_type)
                
                print(f">> [COMPLEX] Slide {slide_idx} analysis:")
                print(f"   - Total objects: {num_objects}")
                print(f"   - Shape types: {list(shape_counts.keys())}")
                print(f"   - Shape counts: {shape_counts}")
                
                # Convert slide to image and process with Gemini
                print(f">> [COMPLEX] Converting slide {slide_idx} to image and processing with Gemini...")
                slide_description = self._convert_slide_to_image_and_process(slide, slide_idx, file_path.name)
                print(f">> [COMPLEX] Slide {slide_idx} description generated: {len(slide_description) if slide_description else 0} characters")
                
                # Format with metadata
                metadata = f"{file_path.name}:slide-{slide_idx}"
                formatted_description = f"## {metadata}\n\n{slide_description}"
                
                complex_slide_descriptions[slide_idx] = formatted_description
                print(f">> [COMPLEX] Slide {slide_idx} processing completed successfully")
            
            print(f">> [COMPLEX] All {len(complex_slides)} complex slides processed")
            return complex_slide_descriptions
            
        except Exception as e:
            print(f">> [COMPLEX] Error processing complex slides: {e}")
            import traceback
            print(f">> [COMPLEX] Full traceback:")
            traceback.print_exc()
            return {}
    
    def _extract_text_from_slide(self, slide) -> str:
        """Extract all text content from a slide"""
        try:
            slide_text_parts = []
            
            # Extract text from all shapes
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text_parts.append(shape.text.strip())
                    
                    # Handle tables specifically to extract cell text
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        table_text = self._extract_text_from_table(shape)
                        if table_text:
                            slide_text_parts.append(table_text)
                    
                    # Handle grouped shapes
                    if hasattr(shape, 'shapes'):
                        for sub_shape in shape.shapes:
                            if hasattr(sub_shape, "text") and sub_shape.text.strip():
                                slide_text_parts.append(sub_shape.text.strip())
                                
                except Exception as shape_error:
                    # Continue processing other shapes if one fails
                    logger.debug(f"Could not extract text from shape: {shape_error}")
                    continue
            
            # Join all text parts with line breaks
            slide_text = "\n".join(slide_text_parts)
            return slide_text
            
        except Exception as e:
            logger.error(f"Error extracting text from slide: {e}")
            return ""
    
    def _extract_text_from_table(self, table_shape) -> str:
        """Extract text content from table shape for text analysis"""
        try:
            table_text_parts = []
            table = table_shape.table
            
            for row in table.rows:
                row_text_parts = []
                for cell in row.cells:
                    cell_text = cell.text.strip() if cell.text else ""
                    if cell_text:
                        row_text_parts.append(cell_text)
                
                if row_text_parts:
                    table_text_parts.append(" | ".join(row_text_parts))
            
            return "\n".join(table_text_parts) if table_text_parts else ""
            
        except Exception as e:
            logger.debug(f"Error extracting text from table: {e}")
            return ""

    def _convert_slide_to_image_and_process(self, slide, slide_idx: int, source_file: str) -> str:
        """Convert a single slide to image and process with Gemini Vision using cross-platform utilities"""
        try:
            print(f">> [IMAGE] Converting slide {slide_idx} to image using cross-platform utilities")
            
            # Extract all text from the slide first
            slide_text = self._extract_text_from_slide(slide)
            print(f">> [IMAGE] Extracted {len(slide_text)} characters of text from slide {slide_idx}")
            
            # Create slide images directory if it doesn't exist
            slide_images_dir = Path("slide_images")
            slide_images_dir.mkdir(exist_ok=True)
            
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            
            # Use the new cross-platform slide exporter
            print(f">> [IMAGE] Attempting cross-platform slide export...")
            image_path = self.slide_exporter.export_slide_as_image(
                self.full_file_path, slide_idx, slide_images_dir, timestamp
            )
            
            if image_path and Path(image_path).exists():
                print(f">> [IMAGE] Slide export successful: {image_path}")
                
                # Process with Gemini Vision including extracted text
                print(f">> [IMAGE] Processing slide {slide_idx} image with Gemini Vision (text + image)")
                description = self._process_slide_image_with_gemini(image_path, slide_idx, source_file, slide_text)
                print(f">> [IMAGE] Generated description for slide {slide_idx}: {len(description) if description else 0} characters")
                
                return description
            else:
                print(f">> [IMAGE] All slide export methods failed, using text-only fallback")
                # Fallback: use extracted text with basic formatting
                if slide_text.strip():
                    print(f">> [IMAGE] Using extracted text as fallback: {len(slide_text)} characters")
                    return f"**Slide {slide_idx} Content (text-only fallback):**\n{slide_text}\n\n*Note: Image processing failed, using text extraction only*"
                else:
                    print(f">> [IMAGE] No text available for fallback")
                    return f"[Slide {slide_idx}: Image processing failed and no text content available]"
            
        except Exception as e:
            print(f">> [IMAGE] Error converting slide {slide_idx} to image: {e}")
            import traceback
            print(f">> [IMAGE] Full traceback:")
            traceback.print_exc()
            return f"[Error processing slide {slide_idx}: {str(e)}]"
    
    def _export_slide_method1_improved(self, slide, slide_idx: int, output_dir: Path, timestamp: str, source_file: str) -> str:
        """Method 1 Improved: Direct slide Export with proper COM handling"""
        powerpoint = None
        presentation = None
        
        try:
            print(f">> [METHOD1] Starting direct slide export for slide {slide_idx}")
            
            # Import COM modules
            import pythoncom
            import comtypes.client
            
            # Initialize COM in apartment threading mode
            pythoncom.CoInitializeEx(pythoncom.COINIT_APARTMENTTHREADED)
            
            # Create PowerPoint application with error handling
            print(f">> [METHOD1] Creating PowerPoint application...")
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            
            # Set visibility carefully
            try:
                powerpoint.Visible = 1  # Use integer instead of boolean
                powerpoint.WindowState = 2  # ppWindowMinimized
            except Exception as vis_error:
                print(f">> [METHOD1] Note: Could not set visibility: {vis_error}")
            
            # Open presentation with full path
            pptx_path = os.path.abspath(source_file)
            print(f">> [METHOD1] Opening: {pptx_path}")
            
            if not os.path.exists(pptx_path):
                raise FileNotFoundError(f"File not found: {pptx_path}")
            
            presentation = powerpoint.Presentations.Open(pptx_path, ReadOnly=1)
            print(f">> [METHOD1] Presentation opened, {presentation.Slides.Count} slides found")
            
            # Validate slide number
            if slide_idx > presentation.Slides.Count:
                raise ValueError(f"Slide {slide_idx} not found (only {presentation.Slides.Count} slides)")
            
            # Get target slide
            target_slide = presentation.Slides[slide_idx]  # PowerPoint uses 1-based indexing
            print(f">> [METHOD1] Target slide {slide_idx} retrieved")
            
            # Create output path with proper Windows formatting
            safe_filename = os.path.basename(source_file).replace('.pptx', '').replace(' ', '_')
            image_filename = f"slide_{slide_idx}_{safe_filename}_{timestamp}.png"
            output_path = output_dir / image_filename
            
            # Ensure the output directory exists and get absolute path
            output_dir.mkdir(parents=True, exist_ok=True)
            abs_output_path = os.path.abspath(str(output_path))
            print(f">> [METHOD1] Output path: {abs_output_path}")
            
            # Export slide with proper parameters
            print(f">> [METHOD1] Exporting slide {slide_idx}...")
            
            # Method 1a: Try direct Export method
            try:
                target_slide.Export(abs_output_path, "PNG", 1920, 1080)
                print(f">> [METHOD1] Direct Export method successful")
                
            except Exception as export_error:
                print(f">> [METHOD1] Direct Export failed: {export_error}")
                print(f">> [METHOD1] Trying alternative export method...")
                
                # Method 1b: Try ExportAsFixedFormat
                temp_dir = f"temp_export_{timestamp}"
                os.makedirs(temp_dir, exist_ok=True)
                
                presentation.ExportAsFixedFormat(
                    Path=os.path.abspath(temp_dir),
                    FixedFormatType=2,  # ppFixedFormatTypePNG
                    Intent=1,  # ppFixedFormatIntentPrint
                    FrameSlides=0,  # Don't frame slides
                    HandoutOrder=1,  # ppPrintHandoutHorizontalFirst
                    OutputType=1,  # ppPrintOutputSlides
                    PrintHiddenSlides=0,  # Don't print hidden slides
                    PrintRange=None,  # Print all slides
                    RangeType=1,  # ppPrintAll
                    SlideShowName="",
                    IncludeDocProps=0,
                    KeepIRMSettings=0,
                    DocStructureTags=0,
                    BitmapMissingFonts=1,
                    UseDocumentICCProfile=0
                )
                
                # Find and copy the specific slide image
                slide_file = os.path.join(temp_dir, f"Slide{slide_idx}.PNG")
                if os.path.exists(slide_file):
                    import shutil
                    shutil.copy2(slide_file, str(output_path))
                    print(f">> [METHOD1] ExportAsFixedFormat method successful")
                    
                    # Cleanup temp directory
                    shutil.rmtree(temp_dir, ignore_errors=True)
                # Find and copy the specific slide image
                slide_file = os.path.join(temp_dir, f"Slide{slide_idx}.PNG")
                if os.path.exists(slide_file):
                    shutil.copy2(slide_file, abs_output_path)
                    print(f">> [METHOD1] ExportAsFixedFormat method successful")
                    
                    # Cleanup temp directory
                    shutil.rmtree(temp_dir, ignore_errors=True)
                else:
                    raise FileNotFoundError(f"Exported slide not found: {slide_file}")
            
            # Verify file was created
            if os.path.exists(abs_output_path):
                file_size = os.path.getsize(abs_output_path)
                print(f">> [METHOD1] SUCCESS: {abs_output_path} ({file_size} bytes)")
                return abs_output_path
            else:
                raise FileNotFoundError("Export completed but file not created")
                
        except Exception as e:
            print(f">> [METHOD1] FAILED: {e}")
            raise e
            
        finally:
            # Clean up COM objects
            try:
                if presentation:
                    presentation.Close()
                    print(f">> [METHOD1] Closed presentation")
            except:
                pass
                
            try:
                if powerpoint:
                    powerpoint.Quit()
                    print(f">> [METHOD1] Quit PowerPoint")
            except:
                pass
                
            try:
                pythoncom.CoUninitialize()
            except:
                pass

    def _export_slide_method2_improved(self, slide_idx: int, source_file: str, output_dir: Path, timestamp: str) -> str:
        """Method 2 Improved: Copy slide with better COM object handling"""
        powerpoint = None
        source_presentation = None
        new_presentation = None
        
        try:
            print(f">> [METHOD2] Starting copy slide method for slide {slide_idx}")
            
            import pythoncom
            import comtypes.client
            
            # Initialize COM
            pythoncom.CoInitializeEx(pythoncom.COINIT_APARTMENTTHREADED)
            
            # Create PowerPoint application
            print(f">> [METHOD2] Creating PowerPoint application...")
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = 1
            
            # Open source presentation
            pptx_path = os.path.abspath(source_file)
            print(f">> [METHOD2] Opening source: {pptx_path}")
            source_presentation = powerpoint.Presentations.Open(pptx_path, ReadOnly=1)
            
            # Create new presentation
            print(f">> [METHOD2] Creating new presentation...")
            new_presentation = powerpoint.Presentations.Add()
            
            # Remove default slide if present
            if new_presentation.Slides.Count > 0:
                new_presentation.Slides[1].Delete()
            
            # Copy target slide
            print(f">> [METHOD2] Copying slide {slide_idx}...")
            source_slide = source_presentation.Slides[slide_idx]
            source_slide.Copy()
            
            # Paste into new presentation
            new_presentation.Slides.Paste()
            print(f">> [METHOD2] Slide copied successfully")
            
            # Export the new presentation with proper path handling
            safe_filename = os.path.basename(source_file).replace('.pptx', '').replace(' ', '_')
            image_filename = f"slide_{slide_idx}_{safe_filename}_{timestamp}.png"
            output_path = output_dir / image_filename
            
            # Ensure the output directory exists and get absolute path
            output_dir.mkdir(parents=True, exist_ok=True)
            abs_output_path = os.path.abspath(str(output_path))
            print(f">> [METHOD2] Output path: {abs_output_path}")
            
            # Get the pasted slide and export it
            pasted_slide = new_presentation.Slides[1]
            pasted_slide.Export(abs_output_path, "PNG", 1920, 1080)
            
            # Verify file creation
            if os.path.exists(abs_output_path):
                file_size = os.path.getsize(abs_output_path)
                print(f">> [METHOD2] SUCCESS: {abs_output_path} ({file_size} bytes)")
                return abs_output_path
            else:
                raise FileNotFoundError("Export completed but file not created")
                
        except Exception as e:
            print(f">> [METHOD2] FAILED: {e}")
            raise e
            
        finally:
            # Cleanup
            try:
                if new_presentation:
                    new_presentation.Close()
            except:
                pass
                
            try:
                if source_presentation:
                    source_presentation.Close()
            except:
                pass
                
            try:
                if powerpoint:
                    powerpoint.Quit()
            except:
                pass
                
            try:
                pythoncom.CoUninitialize()
            except:
                pass

    def _export_slide_method4_libreoffice(self, slide_idx: int, source_file: str, output_dir: Path, timestamp: str) -> str:
        """Method 4: LibreOffice fallback method"""
        try:
            print(f">> [METHOD4] Starting LibreOffice method for slide {slide_idx}")
            
            import subprocess
            
            # First, convert entire PPTX to PDF using LibreOffice
            pdf_path = f"temp_convert_{timestamp}.pdf"
            
            print(f">> [METHOD4] Converting PPTX to PDF with LibreOffice...")
            libreoffice_cmd = [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", ".",
                os.path.abspath(source_file)
            ]
            
            result = subprocess.run(libreoffice_cmd, capture_output=True, text=True, timeout=30)
            
            # LibreOffice creates PDF with same base name
            expected_pdf = os.path.splitext(os.path.basename(source_file))[0] + ".pdf"
            
            if os.path.exists(expected_pdf):
                print(f">> [METHOD4] PDF created: {expected_pdf}")
                
                # Convert PDF page to image using Python
                try:
                    import fitz  # PyMuPDF
                    
                    doc = fitz.open(expected_pdf)
                    if slide_idx <= len(doc):
                        page = doc[slide_idx - 1]  # Convert to 0-based
                        
                        # Render page as image with high quality
                        mat = fitz.Matrix(2.0, 2.0)  # 2x zoom
                        pix = page.get_pixmap(matrix=mat)
                        
                        safe_filename = os.path.basename(source_file).replace('.pptx', '').replace(' ', '_')
                        image_filename = f"slide_{slide_idx}_{safe_filename}_{timestamp}.png"
                        output_path = output_dir / image_filename
                        
                        # Ensure the output directory exists and get absolute path
                        output_dir.mkdir(parents=True, exist_ok=True)
                        abs_output_path = os.path.abspath(str(output_path))
                        pix.save(abs_output_path)
                        
                        doc.close()
                        
                        # Cleanup
                        try:
                            os.remove(expected_pdf)
                        except:
                            pass
                        
                        if os.path.exists(abs_output_path):
                            file_size = os.path.getsize(abs_output_path)
                            print(f">> [METHOD4] SUCCESS: {abs_output_path} ({file_size} bytes)")
                            return abs_output_path
                        
                    else:
                        print(f">> [METHOD4] Slide {slide_idx} not found in PDF (only {len(doc)} pages)")
                        doc.close()
                        
                except ImportError:
                    print(f">> [METHOD4] PyMuPDF not available for PDF to image conversion")
                    
            else:
                print(f">> [METHOD4] PDF conversion failed")
                if result.stderr:
                    print(f">> [METHOD4] LibreOffice error: {result.stderr}")
            
            raise Exception("LibreOffice method failed")
            
        except Exception as e:
            print(f">> [METHOD4] FAILED: {e}")
            raise e
    
    def _export_slide_with_fresh_powerpoint(self, slide_idx: int, source_file: str, output_dir: Path, timestamp: str) -> str:
        """Export slide using a fresh PowerPoint instance (visible application)"""
        powerpoint = None
        presentation = None
        
        try:
            print(f">> [FRESH] Opening fresh PowerPoint instance...")
            import pythoncom
            import comtypes.client
            
            # Initialize COM
            pythoncom.CoInitialize()
            
            # Create PowerPoint application (keep visible to avoid errors)
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = True  # Keep visible
            powerpoint.WindowState = 2  # Minimize but keep visible
            
            # Open the original presentation
            pptx_path = os.path.abspath(source_file)
            print(f">> [FRESH] Opening presentation: {pptx_path}")
            presentation = powerpoint.Presentations.Open(pptx_path, ReadOnly=True)
            
            # Get the specific slide
            target_slide = presentation.Slides[slide_idx]  # PowerPoint uses 1-based indexing
            
            # Create image filename
            image_filename = f"slide_{slide_idx}_{source_file.replace('.pptx', '')}_{timestamp}.png"
            image_path = output_dir / image_filename
            
            print(f">> [FRESH] Exporting slide {slide_idx} to: {image_path}")
            
            # Export slide directly to PNG with high resolution
            target_slide.Export(str(image_path), "PNG", 1920, 1080)
            
            if os.path.exists(str(image_path)):
                file_size = os.path.getsize(str(image_path))
                print(f">> [FRESH] Fresh PowerPoint export successful: {image_path} ({file_size} bytes)")
                return str(image_path)
            else:
                print(f">> [FRESH] Export completed but file not found")
                raise Exception("Export completed but file not created")
                
        except Exception as e:
            print(f">> [FRESH] Fresh PowerPoint export failed: {e}")
            raise e
        
        finally:
            # Cleanup
            try:
                if presentation:
                    presentation.Close()
                    print(f">> [FRESH] Closed presentation")
            except:
                pass
            try:
                if powerpoint:
                    powerpoint.Quit()
                    print(f">> [FRESH] Quit PowerPoint")
            except:
                pass
            try:
                pythoncom.CoUninitialize()
            except:
                pass
    
    def _export_slide_via_pdf_method(self, slide_idx: int, source_file: str, output_dir: Path, timestamp: str) -> str:
        """Export slide by extracting single slide to PDF then converting to image"""
        powerpoint = None
        presentation = None
        single_slide_presentation = None
        
        try:
            print(f">> [PDF] Using PDF method for slide {slide_idx}...")
            import pythoncom
            import comtypes.client
            
            # Initialize COM
            pythoncom.CoInitialize()
            
            # Create PowerPoint application
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = True
            powerpoint.WindowState = 2  # Minimize
            
            # Open original presentation
            pptx_path = os.path.abspath(source_file)
            print(f">> [PDF] Opening presentation: {pptx_path}")
            presentation = powerpoint.Presentations.Open(pptx_path, ReadOnly=True)
            
            # Create new presentation for single slide
            single_slide_presentation = powerpoint.Presentations.Add()
            
            # Remove default slide
            if single_slide_presentation.Slides.Count > 0:
                single_slide_presentation.Slides[1].Delete()
            
            # Copy the target slide
            target_slide = presentation.Slides[slide_idx]
            target_slide.Copy()
            single_slide_presentation.Slides.Paste()
            
            # Create temporary single slide PPTX
            temp_pptx_filename = f"temp_single_slide_{slide_idx}_{timestamp}.pptx"
            temp_pptx_path = output_dir / temp_pptx_filename
            single_slide_presentation.SaveAs(str(temp_pptx_path))
            
            print(f">> [PDF] Created single slide PPTX: {temp_pptx_path}")
            
            # Close PowerPoint objects
            single_slide_presentation.Close()
            presentation.Close()
            powerpoint.Quit()
            pythoncom.CoUninitialize()
            
            # Convert to PDF using LibreOffice
            print(f">> [PDF] Converting single slide to PDF...")
            pdf_path = self._convert_pptx_to_pdf(str(temp_pptx_path))
            
            # Convert PDF to image
            print(f">> [PDF] Converting PDF to image...")
            temp_image_path = self._convert_pdf_to_image(pdf_path)
            
            # Save permanent image
            image_filename = f"slide_{slide_idx}_{source_file.replace('.pptx', '')}_{timestamp}.png"
            permanent_image_path = output_dir / image_filename
            
            import shutil
            shutil.copy2(temp_image_path, permanent_image_path)
            
            # Cleanup temporary files
            try:
                os.unlink(str(temp_pptx_path))
                os.unlink(pdf_path)
                os.unlink(temp_image_path)
            except:
                pass
            
            print(f">> [PDF] PDF method export successful: {permanent_image_path}")
            return str(permanent_image_path)
            
        except Exception as e:
            print(f">> [PDF] PDF method failed: {e}")
            raise e
        
        finally:
            # Ensure cleanup
            try:
                if single_slide_presentation:
                    single_slide_presentation.Close()
            except:
                pass
            try:
                if presentation:
                    presentation.Close()
            except:
                pass
            try:
                if powerpoint:
                    powerpoint.Quit()
            except:
                pass
            try:
                pythoncom.CoUninitialize()
            except:
                pass
    
    def _create_exact_single_slide_copy(self, slide, slide_idx: int, output_dir: Path, timestamp: str, source_file: str) -> str:
        """Create an exact copy of a single slide preserving ALL formatting and content"""
        try:
            from pptx import Presentation as NewPresentation
            import tempfile
            
            # Create new presentation with same slide master/theme
            new_prs = NewPresentation()
            
            # Remove default slide
            if len(new_prs.slides) > 0:
                sp = new_prs.slides._sldIdLst[0]
                new_prs.part.drop_rel(sp.rId)
                del new_prs.slides._sldIdLst[0]
            
            # Get the slide layout (try to preserve original layout)
            try:
                # Use the same slide layout as the original
                original_layout = slide.slide_layout
                # Find a similar layout in the new presentation
                target_layout = new_prs.slide_layouts[0]  # Start with blank layout
                
                # Try to find a layout that matches
                for layout in new_prs.slide_layouts:
                    if layout.name == original_layout.name:
                        target_layout = layout
                        break
                    
            except:
                # Fallback to blank layout
                target_layout = new_prs.slide_layouts[6] if len(new_prs.slide_layouts) > 6 else new_prs.slide_layouts[0]
            
            # Add slide with target layout
            new_slide = new_prs.slides.add_slide(target_layout)
            
            # Copy slide background if possible
            try:
                if hasattr(slide, 'background'):
                    new_slide.background = slide.background
            except:
                pass
            
            print(f">> [EXACT] Creating exact copy of slide {slide_idx} with {len(slide.shapes)} shapes")
            
            # Copy ALL shapes exactly as they are (including charts, images, etc.)
            shapes_copied = 0
            for shape in slide.shapes:
                try:
                    # Try to duplicate the shape exactly
                    self._duplicate_shape_exactly(shape, new_slide)
                    shapes_copied += 1
                except Exception as shape_error:
                    print(f">> [EXACT] Could not copy shape {shapes_copied + 1}: {shape_error}")
                    continue
            
            print(f">> [EXACT] Successfully copied {shapes_copied}/{len(slide.shapes)} shapes")
            
            # Save the exact copy
            exact_copy_filename = f"slide_{slide_idx}_exact_{timestamp}.pptx"
            exact_copy_path = output_dir / exact_copy_filename
            new_prs.save(str(exact_copy_path))
            
            print(f">> [EXACT] Saved exact slide copy: {exact_copy_path}")
            return str(exact_copy_path)
            
        except Exception as e:
            print(f">> [EXACT] Error creating exact slide copy: {e}")
            raise e
    
    def _duplicate_shape_exactly(self, source_shape, target_slide):
        """Duplicate a shape exactly, preserving all properties and content"""
        try:
            # Get basic properties
            left = source_shape.left
            top = source_shape.top
            width = source_shape.width
            height = source_shape.height
            
            # Handle different shape types with better preservation
            shape_type = source_shape.shape_type
            
            if shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                # Text box - preserve text and formatting
                new_shape = target_slide.shapes.add_textbox(left, top, width, height)
                if hasattr(source_shape, 'text'):
                    new_shape.text = source_shape.text
                self._copy_text_formatting(source_shape, new_shape)
                
            elif shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                # Auto shapes (rectangles, circles, etc.) - preserve shape and text
                try:
                    # Try to preserve the exact auto shape type
                    if hasattr(source_shape, 'auto_shape_type'):
                        auto_shape_type = source_shape.auto_shape_type
                        new_shape = target_slide.shapes.add_shape(auto_shape_type, left, top, width, height)
                    else:
                        # Fallback to rectangle
                        new_shape = target_slide.shapes.add_shape(1, left, top, width, height)  # Rectangle
                    
                    if hasattr(source_shape, 'text'):
                        new_shape.text = source_shape.text
                    self._copy_text_formatting(source_shape, new_shape)
                    
                except:
                    # Fallback to text box if auto shape fails
                    new_shape = target_slide.shapes.add_textbox(left, top, width, height)
                    if hasattr(source_shape, 'text'):
                        new_shape.text = source_shape.text
                        
            elif shape_type == MSO_SHAPE_TYPE.TABLE:
                # Tables - try to preserve structure
                if hasattr(source_shape, 'table'):
                    table = source_shape.table
                    rows = len(table.rows)
                    cols = len(table.columns)
                    new_shape = target_slide.shapes.add_table(rows, cols, left, top, width, height)
                    
                    # Copy table content
                    for row_idx in range(rows):
                        for col_idx in range(cols):
                            try:
                                source_cell = table.cell(row_idx, col_idx)
                                target_cell = new_shape.table.cell(row_idx, col_idx)
                                target_cell.text = source_cell.text
                            except:
                                pass
                else:
                    # Fallback representation
                    new_shape = target_slide.shapes.add_textbox(left, top, width, height)
                    new_shape.text = "[TABLE DATA]"
                    
            elif shape_type == MSO_SHAPE_TYPE.CHART:
                # Charts - create placeholder (charts are complex binary objects)
                new_shape = target_slide.shapes.add_textbox(left, top, width, height)
                new_shape.text = "[CHART/DIAGRAM]"
                
            elif shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Pictures - create placeholder (would need image extraction)
                new_shape = target_slide.shapes.add_textbox(left, top, width, height)
                new_shape.text = "[IMAGE]"
                
            elif shape_type == MSO_SHAPE_TYPE.GROUP:
                # Grouped shapes - create placeholder
                new_shape = target_slide.shapes.add_textbox(left, top, width, height)
                new_shape.text = "[GROUPED OBJECTS]"
                
            else:
                # Other shapes - try to preserve text content
                if hasattr(source_shape, 'text') and source_shape.text.strip():
                    new_shape = target_slide.shapes.add_textbox(left, top, width, height)
                    new_shape.text = source_shape.text
                else:
                    # Skip shapes without meaningful content
                    pass
                    
        except Exception as e:
            # If all else fails, create a simple placeholder
            try:
                new_shape = target_slide.shapes.add_textbox(left, top, width, height)
                shape_info = f"[{source_shape.shape_type}]"
                if hasattr(source_shape, 'text') and source_shape.text.strip():
                    shape_info += f": {source_shape.text[:50]}"
                new_shape.text = shape_info
            except:
                pass  # Skip shapes that can't be processed at all
    
    def _copy_text_formatting(self, source_shape, target_shape):
        """Copy text formatting from source to target shape"""
        try:
            if hasattr(source_shape, 'text_frame') and hasattr(target_shape, 'text_frame'):
                source_tf = source_shape.text_frame
                target_tf = target_shape.text_frame
                
                # Copy paragraph-level formatting
                for i, source_para in enumerate(source_tf.paragraphs):
                    if i < len(target_tf.paragraphs):
                        target_para = target_tf.paragraphs[i]
                        
                        # Copy alignment
                        try:
                            target_para.alignment = source_para.alignment
                        except:
                            pass
                            
                        for j, source_run in enumerate(source_para.runs):
                            if j < len(target_para.runs):
                                target_run = target_para.runs[j]
                                try:
                                    if hasattr(source_run, 'font') and hasattr(target_run, 'font'):
                                        target_run.font.size = source_run.font.size
                                        target_run.font.bold = source_run.font.bold
                                        target_run.font.italic = source_run.font.italic
                                except:
                                    pass
        except:
            pass  # Formatting copy is best-effort
    
    def _convert_pdf_to_image(self, pdf_path: str) -> str:
        """Convert PDF to high-quality image"""
        try:
            from PIL import Image
            import fitz  # PyMuPDF
            
            # Open PDF
            doc = fitz.open(pdf_path)
            page = doc[0]  # Get first page
            
            # Convert to image with high DPI for better quality
            mat = fitz.Matrix(2.0, 2.0)  # 2x zoom for better quality
            pix = page.get_pixmap(matrix=mat)
            
            # Save as PNG
            image_path = pdf_path.replace('.pdf', '.png')
            pix.save(image_path)
            
            doc.close()
            return image_path
            
        except ImportError:
            logger.warning("PyMuPDF not available, trying alternative conversion...")
            # Fallback: try with Pillow and pdf2image
            try:
                from pdf2image import convert_from_path
                images = convert_from_path(pdf_path, dpi=200)
                if images:
                    image_path = pdf_path.replace('.pdf', '.png')
                    images[0].save(image_path, 'PNG')
                    return image_path
                else:
                    raise Exception("No images generated from PDF")
            except ImportError:
                raise Exception("Neither PyMuPDF nor pdf2image available for PDF to image conversion")
        except Exception as e:
            logger.error(f"Error converting PDF to image: {e}")
            raise
    
    def _process_slide_image_with_gemini(self, image_path: str, slide_idx: int, source_file: str, slide_text: str = "") -> str:
        """Process slide image with Gemini Vision to extract content, prioritizing extracted text"""
        try:
            # Create enhanced prompt that prioritizes the extracted text
            if slide_text.strip():
                prompt = f"""Analyze this PowerPoint slide image from "{source_file}" (slide {slide_idx}).

EXTRACTED TEXT FROM SLIDE:
{slide_text}

CRITICAL INSTRUCTIONS:
1. PRIORITIZE THE EXTRACTED TEXT ABOVE - This is the PRIMARY SOURCE of information
2. Use the slide image to SUPPLEMENT and ENHANCE the text content, NOT replace it
3. The extracted text contains the exact content from the slide - DO NOT MISS ANY OF IT
4. Generate descriptions that incorporate ALL the text content without omitting anything
5. Use the visual elements in the image to provide context and structure to the text

CRITICAL LANGUAGE PRESERVATION RULE:
- PRESERVE THE ORIGINAL LANGUAGE of all text found in the slide
- If text is in Japanese (日本語), keep it in Japanese - DO NOT translate to English
- If text is in English, keep it in English - DO NOT translate to Japanese
- If slide contains both Japanese and English text, preserve both languages exactly as they appear
- Use the exact same language, characters, and script as found in the original slide and extracted text
- When generating descriptions, match the primary language of the slide content

EXTRACT ONLY THE SEMANTIC BUSINESS INFORMATION - NO VISUAL DESCRIPTIONS.

Focus on extracting information that would be useful for populating a database and answering business queries:

**1. PROCESS FLOW/WORKFLOW**: If this contains a process, workflow, or sequence:
- Use the extracted text as the primary source for step descriptions
- List each step/task in order with its current status (planned, in-progress, completed, blocked, etc.)
- Identify what triggers each step and what outputs it produces
- Extract any dates, deadlines, or timeframes mentioned in the text
- Note dependencies between tasks mentioned in the text
- Identify responsible parties or roles for each step from the text

**2. DATA AND FACTS**: Extract all factual information from the text:
- Include ALL text content exactly as written (titles, labels, data points) from the extracted text
- Any numerical data, percentages, quantities, or metrics mentioned in the text
- Dates, timelines, milestones, and deadlines from the text
- Names of people, roles, departments, or organizations from the text
- Project names, phase names, or task identifiers from the text

**3. STATUS AND PROGRESS**: If the slide shows status information:
- Current state of tasks/projects (use both text and visual indicators)
- Completion percentages or progress indicators from the text
- Issues, risks, or blockers mentioned in the text
- Achievements or completed milestones from the text

DO NOT INCLUDE:
- Colors, shapes, visual design elements
- Layout descriptions or positioning
- Graphical representations (arrows, boxes, charts as visual objects)
- How things "look" or are "displayed"
- Do not try to make up anything that is not present in the extracted text or slides, especially important terms like "project_name", "project_title", etc.

TRANSLATE VISUAL INDICATORS TO MEANING:
- If something appears "completed" (regardless of how it's shown), state "completed"
- If something shows "in progress" (regardless of visual representation), state "in progress"
- Focus on WHAT the visuals represent, not HOW they look
- Use visuals to enhance understanding of the extracted text

CRITICAL: Ensure that EVERY piece of text from the extracted text is incorporated into your description. Do not omit any content from the extracted text.

Provide this as structured business information that could be stored in a database."""
            else:
                # Fallback prompt when no text is extracted
                prompt = f"""Analyze this PowerPoint slide image from "{source_file}" (slide {slide_idx}).

CRITICAL LANGUAGE PRESERVATION RULE:
- PRESERVE THE ORIGINAL LANGUAGE of all text found in the slide
- If text is in Japanese (日本語), keep it in Japanese - DO NOT translate to English
- If text is in English, keep it in English - DO NOT translate to Japanese
- If slide contains both Japanese and English text, preserve both languages exactly as they appear
- Use the exact same language, characters, and script as found in the original slide
- When generating descriptions, match the primary language of the slide content

EXTRACT ONLY THE SEMANTIC BUSINESS INFORMATION - NO VISUAL DESCRIPTIONS.

Focus on extracting information that would be useful for populating a database and answering business queries:

**1. PROCESS FLOW/WORKFLOW**: If this contains a process, workflow, or sequence:
- List each step/task in order with its current status (planned, in-progress, completed, blocked, etc.)
- Identify what triggers each step and what outputs it produces
- Extract any dates, deadlines, or timeframes mentioned
- Note dependencies between tasks
- Identify responsible parties or roles for each step

**2. DATA AND FACTS**: Extract all factual information:
- All text content exactly as written (titles, labels, data points)
- Any numerical data, percentages, quantities, or metrics
- Dates, timelines, milestones, and deadlines
- Names of people, roles, departments, or organizations
- Project names, phase names, or task identifiers

**3. STATUS AND PROGRESS**: If the slide shows status information:
- Current state of tasks/projects (interpret visual indicators as status)
- Completion percentages or progress indicators
- Issues, risks, or blockers mentioned
- Achievements or completed milestones

DO NOT INCLUDE:
- Colors, shapes, visual design elements
- Layout descriptions or positioning
- Graphical representations (arrows, boxes, charts as visual objects)
- How things "look" or are "displayed"
- Do not try to make up anything that is not present on the slides, especially important terms like "project_name", project_title", etc.

TRANSLATE VISUAL INDICATORS TO MEANING:
- If something appears "completed" (regardless of how it's shown), state "completed"
- If something shows "in progress" (regardless of visual representation), state "in progress"
- Focus on WHAT the visuals represent, not HOW they look

Provide this as structured business information that could be stored in a database."""

            # Call Gemini with the slide image
            inputs = [
                ("text", prompt),
                ("file", (image_path, "image/png"))
            ]
            
            response = call_gemini(inputs)
            
            # Extract response text
            if hasattr(response, 'text'):
                description = response.text
            elif hasattr(response, 'candidates') and response.candidates:
                description = response.candidates[0].content.parts[0].text
            else:
                description = str(response)
            
            description = description.strip() if description else ""
            
            if description:
                logger.info(f"✅ Generated description for slide {slide_idx}: {len(description)} characters")
                if slide_text.strip():
                    logger.info(f"✅ Used {len(slide_text)} characters of extracted text as primary source")
                return description
            else:
                logger.warning(f"⚠️ Empty description for slide {slide_idx}")
                return f"[No content extracted from slide {slide_idx}]"
                
        except Exception as e:
            logger.error(f"Error processing slide image with Gemini: {e}")
            return f"[Error processing slide {slide_idx} with Gemini: {str(e)}]"
    
    def _extract_table_data(self, table_shape) -> Optional[List[List[str]]]:
        """Extract data from a PowerPoint table shape"""
        try:
            table_data = []
            table = table_shape.table
            
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text.strip() if cell.text else ""
                    row_data.append(cell_text)
                table_data.append(row_data)
            
            return table_data if table_data else None
            
        except Exception as e:
            logger.error(f"Error extracting table data: {e}")
            return None
    
    def _process_table_with_gemini(self, table_data: List[List[str]], slide_idx: int) -> str:
        """Process table data with Gemini to handle merged cells and formatting"""
        try:
            # Convert table data to CSV format for Gemini
            csv_content = "\n".join([",".join([f'"{cell}"' for cell in row]) for row in table_data])
                    
            prompt = f"""This is a CSV table from slide {slide_idx} that I want you to extract data from. There might be some merged cells, multiple lines in a single cell, etc which I want you to take care of while making the extracted CSV file.

In case of merged cells, unmerge and copy and paste the same content in all the cells.
In case of multiple lines in a single cell, divide them into different rows and copy and paste the other column data for that particular row.
If there is a legend present on the slide, replace the symbols in the csv with the corresponding legend text.
The output should be a well formatted CSV file with all the original data, just the formatting improved and legend text added, if any. Do not make changes just for the sake of making changes, if the CSV is already correctly formatted, return it to me as it is.

Here is the CSV data:
{csv_content}

Please return only the processed CSV data, no additional text or formatting."""

            response = call_gemini([("text", prompt)])

        except Exception as e:
            logger.error(f"Error processing table with Gemini: {e}")
            return ""
            
            # Extract response text
            if hasattr(response, 'text'):
                processed_csv = response.text
            elif hasattr(response, 'candidates') and response.candidates:
                processed_csv = response.candidates[0].content.parts[0].text
            else:
                processed_csv = str(response)
            
            return processed_csv.strip()
            
        except Exception as e:
            logger.error(f"Error processing table with Gemini: {e}")
            # Fallback: return original CSV
            return "\n".join([",".join([f'"{cell}"' for cell in row]) for row in table_data])
    
    def _save_csv_for_pipeline(self, csv_content: str, slide_idx: int, source_file: str) -> str:
        """Save processed CSV to directory for 3-agent pipeline processing"""
        try:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            # Use source file name and slide number for CSV filename
            base_name = Path(source_file).stem
            csv_filename = f"{base_name}_slide{slide_idx}_{timestamp}.csv"
            csv_path = self.csv_output_dir / csv_filename
            
            with open(csv_path, 'w', encoding='utf-8') as f:
                f.write(csv_content)
            
            logger.info(f"💾 Saved CSV for pipeline: {csv_path}")
            return str(csv_path)
            
        except Exception as e:
            logger.error(f"Error saving CSV for pipeline: {e}")
            return ""
    
    def _save_table_csv(self, csv_content: str, slide_idx: int) -> str:
        """Save processed CSV to file"""
        try:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            csv_filename = f"table_slide_{slide_idx}_{timestamp}.csv"
            csv_path = self.table_slides_dir / csv_filename
            
            with open(csv_path, 'w', encoding='utf-8') as f:
                f.write(csv_content)
            
            logger.info(f"💾 Saved table CSV: {csv_path}")
            return str(csv_path)
            
        except Exception as e:
            logger.error(f"Error saving table CSV: {e}")
            return ""
    
    def _save_table_slide(self, slide_content: Dict, source_file: str):
        """Save table slide information to JSON"""
        try:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            json_filename = f"table_slide_{slide_content['slide_number']}_{timestamp}.json"
            json_path = self.table_slides_dir / json_filename
            
            with open(json_path, 'w', encoding='utf-8') as f:
                json.dump(slide_content, f, indent=2, ensure_ascii=False)
            
            logger.info(f"💾 Saved table slide info: {json_path}")
            
        except Exception as e:
            logger.error(f"Error saving table slide info: {e}")
    
    
    def _create_filtered_presentation(self, file_path: Path, slide_numbers: List[int]) -> str:
        """Create a new PPTX with only the specified slides, preserving original slide design"""
        try:
            from pptx import Presentation
            import tempfile
            import shutil
            
            # Load original presentation
            original_prs = Presentation(str(file_path))
            
            # Create timestamp for unique filename
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filtered_path = f"temp_filtered_{timestamp}.pptx"
            
            logger.info(f"📋 Creating new presentation with slides {slide_numbers} from original presentation")
            
            # Create new presentation
            new_prs = Presentation()
            
            # Remove the default blank slide
            if len(new_prs.slides) > 0:
                sp = new_prs.slides._sldIdLst[0]
                new_prs.part.drop_rel(sp.rId)
                del new_prs.slides._sldIdLst[0]
            
            # Copy each specified slide to the new presentation
            for slide_idx in slide_numbers:
                logger.info(f"📄 Copying original slide {slide_idx}")
                
                # Get the source slide (convert to 0-based index)
                source_slide = original_prs.slides[slide_idx - 1]
                
                # Add slide with same layout to new presentation
                try:
                    # Try to use the same slide layout
                    slide_layout = source_slide.slide_layout
                    new_slide = new_prs.slides.add_slide(slide_layout)
                except:
                    # Fallback to blank layout if layout copying fails
                    slide_layout = new_prs.slide_layouts[6]  # Blank layout
                    new_slide = new_prs.slides.add_slide(slide_layout)
                
                # Copy all shapes from source slide
                for shape in source_slide.shapes:
                    try:
                        # Use _copy_shape method or duplicate the shape
                        self._duplicate_shape(shape, new_slide)
                    except Exception as shape_error:
                        logger.warning(f"Could not copy shape {getattr(shape, 'name', 'unknown')}: {shape_error}")
                        continue
            
            # Save the new presentation
            new_prs.save(filtered_path)
            
            logger.info(f"✅ Created filtered PPTX with {len(slide_numbers)} slides: {filtered_path}")
            return filtered_path
            
        except Exception as e:
            logger.error(f"Error creating filtered presentation: {e}")
            # Fallback: create timestamp-based filename and return empty
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            fallback_path = f"temp_filtered_{timestamp}.pptx"
            return fallback_path
    
    def _duplicate_shape(self, source_shape, target_slide):
        """Duplicate a shape from source to target slide"""
        try:
            # Get shape properties
            left = source_shape.left
            top = source_shape.top
            width = source_shape.width
            height = source_shape.height
            
            # Handle different shape types
            if source_shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                # Text box
                new_shape = target_slide.shapes.add_textbox(left, top, width, height)
                if hasattr(source_shape, 'text'):
                    new_shape.text = source_shape.text
                    
            elif source_shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                # Auto shape (rectangles, etc.)
                if hasattr(source_shape, 'text'):
                    new_shape = target_slide.shapes.add_textbox(left, top, width, height)
                    new_shape.text = source_shape.text
                    
            elif source_shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                # Skip placeholders as they're part of the layout
                pass
                
            elif source_shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                # Tables - simplified copy
                if hasattr(source_shape, 'text'):
                    new_shape = target_slide.shapes.add_textbox(left, top, width, height)
                    new_shape.text = f"[TABLE: {source_shape.text[:100]}...]"
                    
            else:
                # Other shapes - try basic text copying
                if hasattr(source_shape, 'text') and source_shape.text.strip():
                    new_shape = target_slide.shapes.add_textbox(left, top, width, height)
                    new_shape.text = source_shape.text
                    
        except Exception as e:
            # Skip shapes that can't be copied
            pass
    
    def _convert_pptx_to_pdf(self, pptx_path: str) -> str:
        """Convert PPTX to PDF using cross-platform office utilities"""
        
        print(f">> [CONVERT] Starting PPTX to PDF conversion for: {os.path.basename(pptx_path)}")
        print(f">> [CONVERT] Input file exists: {os.path.exists(pptx_path)}")
        print(f">> [CONVERT] Input file size: {os.path.getsize(pptx_path) if os.path.exists(pptx_path) else 'N/A'} bytes")
        
        try:
            # Use cross-platform office utilities (Windows COM first, LibreOffice fallback)
            output_dir = Path(pptx_path).parent
            pdf_filename = f"{Path(pptx_path).stem}.pdf"
            pdf_path = output_dir / pdf_filename
            
            print(f">> [CONVERT] Using cross-platform conversion utilities...")
            result_path = ppt_to_pdf(pptx_path, pdf_path)
            print(f">> [CONVERT] Conversion successful: {result_path}")
            return result_path
            
        except OfficeConversionError as e:
            print(f">> [CONVERT] Cross-platform conversion failed: {e}")
            raise Exception(f"PowerPoint to PDF conversion failed: {e}")
    
    def _convert_with_com(self, pptx_path: str) -> str:
        """Convert PPTX to PDF using PowerPoint COM automation with simplified error handling"""
        powerpoint = None
        presentation = None
        
        try:
            print(f">> [COM] Initializing COM for: {os.path.basename(pptx_path)}")
            # Initialize COM
            pythoncom.CoInitialize()
            
            # Create PowerPoint application
            try:
                print(f">> [COM] Creating PowerPoint COM object...")
                powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
                print(f">> [COM] PowerPoint COM object created successfully")
            except Exception as com_error:
                print(f">> [COM] Failed to create PowerPoint COM object: {com_error}")
                raise Exception(f"Failed to create PowerPoint COM object: {com_error}")
            
            # Check if PowerPoint object was created properly
            if not powerpoint:
                print(f">> [COM] PowerPoint COM object is None")
                raise Exception("PowerPoint COM object is None")
            
            # Try to access Presentations - this is where the error occurs
            try:
                print(f">> [COM] Accessing Presentations collection...")
                presentations = powerpoint.Presentations
                print(f">> [COM] Presentations collection accessed successfully")
            except Exception as pres_error:
                print(f">> [COM] Failed to access Presentations collection: {pres_error}")
                raise Exception(f"Failed to access Presentations collection: {pres_error}")
            
            powerpoint.Visible = True  # Keep visible to avoid errors
            print(f">> [COM] Set PowerPoint visible")
            
            # Open presentation with absolute path
            abs_pptx_path = os.path.abspath(pptx_path)
            print(f">> [COM] Opening presentation: {abs_pptx_path}")
            
            try:
                presentation = presentations.Open(abs_pptx_path, ReadOnly=True)
                print(f">> [COM] Presentation opened successfully")
            except Exception as open_error:
                print(f">> [COM] Failed to open presentation: {open_error}")
                raise Exception(f"Failed to open presentation: {open_error}")
            
            # Export to PDF
            pdf_path = pptx_path.replace('.pptx', '.pdf')
            abs_pdf_path = os.path.abspath(pdf_path)
            print(f">> [COM] Exporting to PDF: {abs_pdf_path}")
            
            try:
                # Use ExportAsFixedFormat for better control
                print(f">> [COM] Attempting ExportAsFixedFormat method...")
                presentation.ExportAsFixedFormat(abs_pdf_path, 2)  # 2 = ppFixedFormatTypePDF
                print(f">> [COM] ExportAsFixedFormat successful")
                print(f">> [COM] Checking if PDF was created: {os.path.exists(pdf_path)}")
                if os.path.exists(pdf_path):
                    print(f">> [COM] PDF file size: {os.path.getsize(pdf_path)} bytes")
                    print(f">> [COM] COM conversion successful: {pdf_path}")
                    return pdf_path
                else:
                    print(f">> [COM] PDF file was not created despite no errors")
                    raise Exception("PDF file was not created")
                    
            except Exception as export_error:
                print(f">> [COM] ExportAsFixedFormat failed: {export_error}")
                # Try alternative method: SaveAs with ppSaveAsPDF
                try:
                    print(f">> [COM] Attempting SaveAs method...")
                    presentation.SaveAs(abs_pdf_path, 32)  # 32 = ppSaveAsPDF
                    print(f">> [COM] SaveAs method completed")
                    if os.path.exists(pdf_path):
                        print(f">> [COM] COM conversion successful (SaveAs method): {pdf_path}")
                        return pdf_path
                    else:
                        print(f">> [COM] SaveAs method did not create PDF file")
                        raise Exception("SaveAs method did not create PDF file")
                except Exception as save_error:
                    print(f">> [COM] SaveAs method failed: {save_error}")
                    raise Exception(f"Both ExportAsFixedFormat and SaveAs failed: {export_error}, {save_error}")
            
        except Exception as e:
            print(f">> [COM] COM automation failed: {e}")
            raise Exception(f"PowerPoint COM automation failed: {e}")
        
        finally:
            # Ensure proper cleanup
            print(f">> [COM] Cleaning up COM objects...")
            try:
                if presentation:
                    presentation.Close()
                    print(f">> [COM] Presentation closed")
            except:
                print(f">> [COM] Error closing presentation")
                pass
            try:
                if powerpoint:
                    powerpoint.Quit()
                    print(f">> [COM] PowerPoint quit")
            except:
                print(f">> [COM] Error quitting PowerPoint")
                pass
            try:
                pythoncom.CoUninitialize()
                print(f">> [COM] COM uninitialized")
            except:
                print(f">> [COM] Error uninitializing COM")
                pass
    
    def _convert_with_libreoffice(self, pptx_path: str) -> str:
        """Fallback: Convert PPTX to PDF using LibreOffice"""
        try:
            pdf_path = pptx_path.replace('.pptx', '.pdf')
            output_dir = os.path.dirname(os.path.abspath(pptx_path))
            abs_pptx_path = os.path.abspath(pptx_path)
            
            print(f">> [LIBREOFFICE] Converting: {abs_pptx_path}")
            print(f">> [LIBREOFFICE] Output dir: {output_dir}")
            print(f">> [LIBREOFFICE] Expected PDF: {pdf_path}")
            
            # Try multiple LibreOffice command variations (verified path first)
            commands_to_try = [
                [r'C:\Program Files\LibreOffice\program\soffice.exe', '--headless', '--convert-to', 'pdf', '--outdir', output_dir, abs_pptx_path],
                ['soffice', '--headless', '--convert-to', 'pdf', '--outdir', output_dir, abs_pptx_path],
                ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', output_dir, abs_pptx_path],
                ['C:/Program Files/LibreOffice/program/soffice.exe', '--headless', '--convert-to', 'pdf', '--outdir', output_dir, abs_pptx_path]
            ]
            
            for i, cmd in enumerate(commands_to_try, 1):
                try:
                    print(f">> [LIBREOFFICE] Attempt {i}: {' '.join(cmd)}")
                    result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
                    
                    print(f">> [LIBREOFFICE] Return code: {result.returncode}")
                    if result.stdout:
                        print(f">> [LIBREOFFICE] STDOUT: {result.stdout}")
                    if result.stderr:
                        print(f">> [LIBREOFFICE] STDERR: {result.stderr}")
                    
                    if result.returncode == 0 and os.path.exists(pdf_path):
                        print(f">> [LIBREOFFICE] Conversion successful: {pdf_path}")
                        print(f">> [LIBREOFFICE] PDF file size: {os.path.getsize(pdf_path)} bytes")
                        return pdf_path
                    elif result.stderr:
                        print(f">> [LIBREOFFICE] Attempt {i} failed with stderr: {result.stderr}")
                        
                except FileNotFoundError:
                    print(f">> [LIBREOFFICE] Command not found: {cmd[0]}")
                    continue
                except subprocess.TimeoutExpired:
                    print(f">> [LIBREOFFICE] Attempt {i} timed out after 120 seconds")
                    continue
                except Exception as e:
                    print(f">> [LIBREOFFICE] Attempt {i} error: {e}")
                    continue
            
            print(f">> [LIBREOFFICE] All conversion attempts failed")
            raise Exception("All LibreOffice conversion attempts failed")
                
        except Exception as e:
            print(f">> [LIBREOFFICE] LibreOffice conversion error: {e}")
            raise
    
    def _process_images_in_markdown(self, markdown_content: str, source_file: str) -> str:
        """Extract images from markdown and replace with Gemini descriptions, including context text"""
        if not markdown_content:
            return markdown_content
        
        try:
            # Find all image references in markdown - multiple patterns
            patterns = [
                r'!\[([^\]]*)\]\(data:image/[^;]+;base64,([^)]+)\)',  # Base64 images
                r'!\[([^\]]*)\]\(([^)]+\.(?:png|jpg|jpeg|gif|svg))\)',  # Image file links
                r'\[Image:\s*([^\]]+)\]',  # Simple [Image: description] format
                r'<img[^>]+src=["\']([^"\']+)["\'][^>]*>',  # HTML img tags
            ]
            
            all_images = []
            processed_markdown = markdown_content
            
            for pattern in patterns:
                images = re.findall(pattern, markdown_content, re.IGNORECASE)
                all_images.extend(images)
            
            logger.info(f"🖼️ Found {len(all_images)} images in markdown to process")
            
            # If no images found, log the first 500 chars to debug
            if not all_images:
                logger.info(f"🔍 No images found. Sample markdown content: {markdown_content[:500]}...")
            
            for img_data in all_images:
                if isinstance(img_data, tuple) and len(img_data) >= 2:
                    img_name, base64_or_url = img_data[0], img_data[1]
                else:
                    img_name, base64_or_url = str(img_data), ""
                
                # Extract surrounding context text for better understanding
                context_text = self._extract_context_around_image(markdown_content, img_name)
                
                # Generate description with Gemini
                description = self._generate_image_description_with_gemini(base64_or_url, img_name, source_file, context_text)
                
                # Skip replacement if it's a logo/decorative image (empty description)
                if not description.strip():
                    logger.info(f"🏷️ Skipping replacement for logo/decorative image: {img_name}")
                    continue
                
                # Replace image with description - handle base64 data safely
                # Find the exact image markdown and replace it
                img_markdown = f"![{img_name}](data:image/"
                start_pos = processed_markdown.find(img_markdown)
                
                if start_pos != -1:
                    # Find the end of the image markdown
                    end_pos = processed_markdown.find(")", start_pos)
                    if end_pos != -1:
                        # Extract the full image markdown
                        full_img_markdown = processed_markdown[start_pos:end_pos + 1]
                        
                        # Replace with description
                        new_text = f"\n\n**[DIAGRAM: {img_name}]**\n{description}\n"
                        processed_markdown = processed_markdown.replace(full_img_markdown, new_text)
                        
                        logger.info(f"✅ Replaced image {img_name} with Gemini description")
                    else:
                        logger.warning(f"⚠️ Could not find end of image markdown for {img_name}")
                else:
                    logger.warning(f"⚠️ Could not find image markdown for {img_name}")
            
            return processed_markdown
            
        except Exception as e:
            logger.error(f"Error processing images in markdown: {e}")
            return markdown_content
    
    def _extract_context_around_image(self, markdown_content: str, img_name: str) -> str:
        """Extract surrounding context text around an image for better understanding"""
        try:
            # Find the image in the markdown
            img_pattern = f"!\\[{re.escape(img_name)}\\]"
            match = re.search(img_pattern, markdown_content)
            
            if not match:
                return ""
            
            # Get position of the image
            img_pos = match.start()
            
            # Extract context before and after the image (up to 500 characters each direction)
            start_pos = max(0, img_pos - 500)
            end_pos = min(len(markdown_content), img_pos + 500)
            
            context_before = markdown_content[start_pos:img_pos].strip()
            context_after = markdown_content[match.end():end_pos].strip()
            
            # Clean up the context (remove other image references, excessive whitespace)
            context_parts = []
            if context_before:
                # Get last few lines before the image
                context_before_lines = context_before.split('\n')[-3:]
                context_before_clean = '\n'.join(context_before_lines).strip()
                if context_before_clean:
                    context_parts.append("BEFORE: " + context_before_clean)
            
            if context_after:
                # Get first few lines after the image
                context_after_lines = context_after.split('\n')[:3]
                context_after_clean = '\n'.join(context_after_lines).strip()
                if context_after_clean:
                    context_parts.append("AFTER: " + context_after_clean)
            
            context_text = '\n'.join(context_parts)
            
            if context_text:
                logger.debug(f"Extracted {len(context_text)} characters of context for image {img_name}")
            
            return context_text
            
        except Exception as e:
            logger.debug(f"Could not extract context for image {img_name}: {e}")
            return ""
    
    def _generate_image_description_with_gemini(self, base64_data: str, img_name: str, source_file: str, context_text: str = "") -> str:
        """Generate description for an image using Gemini Vision, optionally with surrounding context text"""
        try:
            # Create prompt for image analysis with context text if available
            if context_text.strip():
                prompt = f"""Analyze this image from PowerPoint file "{source_file}" (image: {img_name}).

SURROUNDING CONTEXT TEXT:
{context_text}

CRITICAL INSTRUCTIONS:
1. Use the surrounding context text to better understand what this image represents
2. Focus on the business/technical content, not visual design
3. The context text provides important information about the slide content

CRITICAL LANGUAGE PRESERVATION RULE:
- PRESERVE THE ORIGINAL LANGUAGE of all text found in the image and context
- If text is in Japanese (日本語), keep it in Japanese - DO NOT translate to English
- If text is in English, keep it in English - DO NOT translate to Japanese
- If image contains both Japanese and English text, preserve both languages exactly as they appear
- Use the exact same language, characters, and script as found in the original image and context
- When generating descriptions, match the primary language of the image content

EXTRACT ONLY THE SEMANTIC BUSINESS INFORMATION - NO VISUAL DESCRIPTIONS.

Focus on extracting information useful for database population and business queries:

**1. PROCESS FLOW/WORKFLOW**: If this shows a process or workflow:
- List each step/task with current status (planned, in-progress, completed, etc.)
- Extract timelines, deadlines, or dates mentioned
- Note dependencies and sequence between tasks
- Identify responsible parties or roles

**2. BUSINESS DATA**: Extract all factual information:
- All text content exactly as written
- Names, titles, roles, departments
- Project names, phase names, identifiers
- Numerical data, metrics, percentages
- Dates and milestones

**3. STATUS INFORMATION**: Interpret visual indicators as status:
- Task completion status (completed, in-progress, planned, blocked)
- Progress indicators or percentages
- Issues, risks, or achievements mentioned

DO NOT INCLUDE:
- Visual design elements (colors, shapes, positioning)
- How things "look" or "appear"
- Layout or formatting descriptions

TRANSLATE VISUAL INDICATORS TO MEANING:
- Focus on WHAT visual elements represent, not HOW they look
- Convert visual status indicators to semantic status information

Provide structured business information suitable for database storage."""
            else:
                prompt = f"""Analyze this image from PowerPoint file "{source_file}" (image: {img_name}).

CRITICAL LANGUAGE PRESERVATION RULE:
- PRESERVE THE ORIGINAL LANGUAGE of all text found in the image
- If text is in Japanese (日本語), keep it in Japanese - DO NOT translate to English
- If text is in English, keep it in English - DO NOT translate to Japanese
- If image contains both Japanese and English text, preserve both languages exactly as they appear
- Use the exact same language, characters, and script as found in the original image
- When generating descriptions, match the primary language of the image content

EXTRACT ONLY THE SEMANTIC BUSINESS INFORMATION - NO VISUAL DESCRIPTIONS.

Focus on extracting information useful for database population and business queries:

**1. PROCESS FLOW/WORKFLOW**: If this shows a process or workflow:
- List each step/task with current status (planned, in-progress, completed, etc.)
- Extract timelines, deadlines, or dates mentioned
- Note dependencies and sequence between tasks
- Identify responsible parties or roles

**2. BUSINESS DATA**: Extract all factual information:
- All text content exactly as written
- Names, titles, roles, departments
- Project names, phase names, identifiers
- Numerical data, metrics, percentages
- Dates and milestones

**3. STATUS INFORMATION**: Interpret visual indicators as status:
- Task completion status (completed, in-progress, planned, blocked)
- Progress indicators or percentages
- Issues, risks, or achievements mentioned

DO NOT INCLUDE:
- Visual design elements (colors, shapes, positioning)
- How things "look" or "appear"
- Layout or formatting descriptions

TRANSLATE VISUAL INDICATORS TO MEANING:
- Focus on WHAT visual elements represent, not HOW they look
- Convert visual status indicators to semantic status information

Provide structured business information suitable for database storage."""

            # Decode base64 to bytes for Gemini
            image_bytes = base64.b64decode(base64_data)
            
            # Create temporary image file
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_file:
                temp_file.write(image_bytes)
                temp_path = temp_file.name
            
            try:
                # Call Gemini with image
                inputs = [
                    ("text", prompt),
                    ("file", (temp_path, "image/png"))
                ]
                
                response = call_gemini(inputs)
                
                # Extract response text
                if hasattr(response, 'text'):
                    description = response.text
                elif hasattr(response, 'candidates') and response.candidates:
                    description = response.candidates[0].content.parts[0].text
                else:
                    description = str(response)
                
                # Clean up the description
                description = description.strip() if description else ""
                
                # Handle logo/decorative element detection
                if not description or description.lower() in ['""', '""', 'empty', 'logo', 'decorative']:
                    logger.info(f"🏷️ Detected logo/decorative image {img_name} - skipping description")
                    return ""
                
                logger.info(f"✅ Generated description for image {img_name}: {len(description)} characters")
                if context_text.strip():
                    logger.info(f"✅ Used {len(context_text)} characters of context text")
                return description
                
            finally:
                # Cleanup temp file
                try:
                    os.unlink(temp_path)
                except:
                    pass
            
        except Exception as e:
            logger.error(f"Error generating image description with Gemini: {e}")
            return f"[Description unavailable for {img_name} - {str(e)}]"
    
    def _extract_fallback_text_from_slides(self, source_file: str, simple_slides: List[int]) -> str:
        """Extract basic text content from simple slides using direct text extraction"""
        try:
            from pptx import Presentation
            
            # Load the presentation
            prs = Presentation(source_file)
            
            fallback_content = []
            
            for slide_idx in simple_slides:
                try:
                    slide = prs.slides[slide_idx - 1]  # Convert to 0-based index
                    
                    # Extract all text from the slide
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    
                    if slide_text:
                        slide_content = f"## {Path(source_file).name}:slide-{slide_idx}\n\n"
                        slide_content += "\n\n".join(slide_text)
                        fallback_content.append(slide_content)
                        logger.info(f"✅ Extracted fallback text from slide {slide_idx}: {len(' '.join(slide_text))} characters")
                    else:
                        logger.warning(f"⚠️ No text found in slide {slide_idx}")
                        
                except Exception as slide_error:
                    logger.error(f"Error extracting text from slide {slide_idx}: {slide_error}")
                    continue
            
            result = "\n\n".join(fallback_content)
            logger.info(f"💾 Generated fallback content for {len(simple_slides)} slides: {len(result)} characters total")
            return result
            
        except Exception as e:
            logger.error(f"Error extracting fallback text from slides: {e}")
            return ""
    
    def _process_csvs_and_combine_results(self, markdown_content: str, processed_tables: Dict, source_file: str, raw_ocr_path: Path = None, pdf_file_path: str = None, complex_slide_descriptions: Dict[int, str] = None, simple_slides: List[int] = None) -> Dict[str, Any]:
        """Process CSVs through 3-agent pipeline and combine all results"""
        try:
            current_time = datetime.now().isoformat()
            
            # Create enhanced markdown by combining text extraction content with complex slide descriptions and table descriptions
            enhanced_markdown = ""
            
            # Add text extraction content (simple slides with metadata)
            if markdown_content:
                enhanced_markdown += f"# Simple Slides (Text Extraction)\n\n"
                enhanced_markdown += markdown_content
                enhanced_markdown += "\n\n"
            
            # Add complex slide descriptions with metadata (already formatted)
            if complex_slide_descriptions:
                enhanced_markdown += f"# Complex Slides (Gemini Vision Analysis)\n\n"
                for slide_idx in sorted(complex_slide_descriptions.keys()):
                    enhanced_markdown += complex_slide_descriptions[slide_idx]
                    enhanced_markdown += "\n\n"
            
            # Add table slide descriptions (processed with Gemini)
            if processed_tables:
                enhanced_markdown += f"# Table Slides (Gemini Analysis)\n\n"
                for slide_idx in sorted(processed_tables.keys()):
                    table_info = processed_tables[slide_idx]
                    enhanced_markdown += f"## {source_file}:slide-{slide_idx}\n\n"
                    enhanced_markdown += table_info.get('description', '')
                    enhanced_markdown += "\n\n"
            
            # Save the final enhanced markdown file
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            final_markdown_filename = f"{Path(source_file).stem}_final_enhanced_{timestamp}.md"
            final_markdown_path = Path(final_markdown_filename)

            with open(final_markdown_path, 'w', encoding='utf-8') as f:
                f.write(f"# Complete Analysis of {source_file}\n\n")
                f.write(f"*Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*\n\n")
                f.write(enhanced_markdown)
            
            logger.info(f"💾 Saved final enhanced markdown: {final_markdown_filename}")
            
            # Process the enhanced markdown file through 3-agent pipeline
            enhanced_md_pipeline_result = None
            try:
                logger.info(f"🔄 Processing enhanced markdown through 3-agent pipeline...")
                enhanced_md_pipeline_result = self._run_3_agent_pipeline_for_markdown(str(final_markdown_path), source_file)
                logger.info(f"✅ Enhanced markdown processed through 3-agent pipeline successfully")
            except Exception as e:
                logger.error(f"❌ Failed to process enhanced markdown through 3-agent pipeline: {e}")
                enhanced_md_pipeline_result = {'error': str(e)}
            
            # Process markdown content into structured data with proper source attribution
            markdown_data = []
            if markdown_content:
                sections = self._extract_sections_from_markdown(markdown_content, source_file)
                markdown_data.extend(sections)
            
            # Add complex slide content to markdown data
            if complex_slide_descriptions:
                for slide_idx, description in complex_slide_descriptions.items():
                    markdown_data.append({
                        "source": f"{source_file}:slide-{slide_idx}",
                        "title": f"Complex Slide {slide_idx} Analysis",
                        "content": description,
                        "last_updated_date": current_time
                    })

            # Add table slide content to markdown data (simplified structure)
            table_processing_results = []
            if processed_tables:
                for slide_idx, table_info in processed_tables.items():
                    # Add table slide description to markdown data
                    markdown_data.append({
                        "source": f"{source_file}:slide-{slide_idx}",
                        "title": f"Table Slide {slide_idx} Analysis",
                        "content": table_info.get('description', ''),
                        "last_updated_date": current_time
                    })
                    
                    # Track table processing (no CSV pipeline for tables now)
                    table_processing_results.append({
                        'slide_number': slide_idx,
                        'processing_method': 'gemini_vision_only',
                        'source': f"{source_file}:slide-{slide_idx}",
                        'description_length': len(table_info.get('description', '')),
                        'extracted_text_length': len(table_info.get('extracted_text', ''))
                    })
            
            # Combine all results
            return {
                "file_type": "powerpoint_text_extraction",
                "source_file": source_file,
                "processing_method": "text_extraction_with_complexity_filtering_and_gemini_only_tables",
                "total_markdown_entries": len(markdown_data),
                "total_table_slides_processed": len(table_processing_results),
                "total_complex_slides_identified": len(complex_slide_descriptions) if complex_slide_descriptions else 0,
                "last_updated_date": current_time,
                "raw_ocr_markdown_file": str(raw_ocr_path) if raw_ocr_path else None,
                "pdf_file_path": pdf_file_path if pdf_file_path else None,
                "final_enhanced_markdown_file": str(final_markdown_path),
                "enhanced_markdown": enhanced_markdown,  # Include the enhanced markdown with complex slides and tables
                "enhanced_markdown_pipeline_result": enhanced_md_pipeline_result,  # 3-agent pipeline result for enhanced markdown
                "markdown_data": markdown_data,
                "table_processing_results": table_processing_results,  # Simplified table results
                "complex_slides_info": complex_slide_descriptions if complex_slide_descriptions else {},
                "table_slides_info": processed_tables if processed_tables else {},
                "summary": {
                    "markdown_entries": len(markdown_data),
                    "table_slides_processed_with_gemini": len(table_processing_results),
                    "complex_slides_identified": len(complex_slide_descriptions) if complex_slide_descriptions else 0,
                    "complex_slides_processed": len(complex_slide_descriptions) if complex_slide_descriptions else 0,
                    "enhanced_markdown_pipeline_success": enhanced_md_pipeline_result.get('status') == 'success' if enhanced_md_pipeline_result else False
                }
            }
            
        except Exception as e:
            logger.error(f"Error processing CSVs and combining results: {e}")
            return {
                "file_type": "powerpoint_text_extraction",
                "source_file": source_file,
                "error": str(e),
                "last_updated_date": datetime.now().isoformat(),
                "markdown_data": [],
                "table_processing_results": []
            }
    
    def _extract_sections_from_markdown(self, markdown_content: str, source_file: str) -> List[Dict]:
        """Extract sections from markdown with proper slide attribution"""
        sections = []
        current_time = datetime.now().isoformat()
        
        # Try to identify which slides the content came from
        # This is based on the text extraction metadata
        pages = markdown_content.split('\n\n')
        
        for i, section in enumerate(pages):
            section = section.strip()
            if section and len(section) > 50:  # Only substantial content
                title = self._extract_title_from_section(section)
                
                # Try to determine slide number from content or use sequential numbering
                slide_number = i + 1  # Approximate slide attribution
                
                sections.append({
                    "source": f"{source_file}:slide-{slide_number}",
                    "title": title,
                    "content": section,
                    "last_updated_date": current_time
                })
        
        return sections
    
    def _run_3_agent_pipeline_for_csv(self, csv_path: str, source_file: str, slide_idx: int) -> Dict:
        """Run the 3-agent pipeline for a specific CSV file"""
        try:
            import subprocess
            import sys
            
            # Run the 3-agent pipeline script
            cmd = [
                sys.executable, "run_3_agents.py", 
                csv_path  # Database insertion enabled (no --no-insert flag)
            ]
            
            logger.info(f"Running command: {' '.join(cmd)}")
            result = subprocess.run(cmd, capture_output=True, text=True, cwd=os.getcwd())
            
            if result.returncode == 0:
                logger.info(f"✅ 3-agent pipeline completed for {csv_path}")
                return {
                    'status': 'success',
                    'stdout': result.stdout,
                    'processing_completed': True
                }
            else:
                logger.error(f"❌ 3-agent pipeline failed for {csv_path}: {result.stderr}")
                return {
                    'status': 'failed',
                    'stderr': result.stderr,
                    'stdout': result.stdout,
                    'processing_completed': False
                }
                
        except Exception as e:
            logger.error(f"Exception running 3-agent pipeline for {csv_path}: {e}")
            return {
                'status': 'exception',
                'error': str(e),
                'processing_completed': False
            }
    
    def _run_3_agent_pipeline_for_markdown(self, markdown_path: str, source_file: str) -> Dict:
        """Run the modified 3-agent pipeline for PowerPoint-generated markdown (skip Agent 2)"""
        try:
            import subprocess
            import sys
            
            print(f">> [PPT-PIPELINE] Starting modified 3-agent pipeline for PowerPoint markdown...")
            print(f">> [PPT-PIPELINE] Markdown file: {os.path.basename(markdown_path)}")
            print(f">> [PPT-PIPELINE] Source PowerPoint: {source_file}")
            
            # Generate unique output filenames with timestamp
            timestamp = datetime.now().strftime("%H%M%S")
            base_name = f"{Path(source_file).stem}_ppt_pipeline"
            agent1_output = f"{base_name}_{timestamp}_agent1_output.json"
            agent3_output = f"{base_name}_{timestamp}_agent3_output.json"
            
            # Step 1: Agent 1 - Raw data extraction for markdown
            print(f">> [PPT-PIPELINE] AGENT 1: Raw data extraction from markdown...")
            try:
                import data_extraction as agent1
                
                # Load schema and attribute descriptions
                schema = agent1.load_schema("schema_postgresql.json")
                attribute_descriptions = agent1.read_attribute_descriptions_csv("attributes_description.csv")
                
                # Process the markdown file
                agent1.process_text_file(markdown_path, schema, attribute_descriptions, "pydantic_models.py", agent1_output)
                print(f">> [PPT-PIPELINE] Agent 1 completed: {agent1_output}")
                
            except Exception as e:
                print(f">> [PPT-PIPELINE] Agent 1 failed: {e}")
                return {
                    'status': 'failed',
                    'error': f"Agent 1 failed: {e}",
                    'agent_failed': 1,
                    'processing_completed': False
                }
            
            # Step 2: Direct Database Insertion (Skip Agent 2 validation)
            print(f">> [PPT-PIPELINE] DIRECT INSERT: Inserting Agent 1 data into PostgreSQL (skipping validation)...")
            insertion_errors = []
            insertion_result = None
            error_log_file = None
            
            try:
                from insert_to_postgresql import PostgreSQLInserter
                inserter = PostgreSQLInserter()
                
                # Insert Agent 1 data directly and capture any errors
                insertion_result = inserter.insert_json_file(agent1_output)
                
                # Check for insertion errors
                insertion_errors = insertion_result.get('errors', [])
                inserted_records = insertion_result.get('inserted_records', 0)
                skipped_records = insertion_result.get('skipped_records', 0)
                
                print(f">> [PPT-PIPELINE] Direct insertion results:")
                if inserted_records > 0:
                    print(f"   - {inserted_records} records inserted successfully")
                if skipped_records > 0:
                    print(f"   - {skipped_records} records had issues")
                if insertion_errors:
                    print(f"   - {len(insertion_errors)} insertion errors found")
                    
                    # Look for the error log file
                    expected_error_file = f"agent1_insertion_errors_{os.path.basename(agent1_output)}"
                    if os.path.exists(expected_error_file):
                        error_log_file = expected_error_file
                        print(f"   - Error log created: {error_log_file}")
                else:
                    print(f"   - No insertion errors")
                    
            except Exception as e:
                print(f">> [PPT-PIPELINE] Direct database insertion failed: {e}")
                insertion_errors.append(f"Database insertion exception: {e}")
            
            # Step 3: Agent 3 - Error rectification (only if there were errors)
            if insertion_errors or error_log_file:
                print(f">> [PPT-PIPELINE] AGENT 3: Error rectification and fixing...")
                try:
                    from agent3_relationship_fixer import Agent3RelationshipFixer
                    agent3 = Agent3RelationshipFixer()
                    
                    # Process errors with Agent 3
                    if error_log_file and os.path.exists(error_log_file):
                        print(f">> [PPT-PIPELINE] Agent 3 processing errors from: {error_log_file}")
                        fixed_data = agent3.process_data(error_log_file, agent3_output, error_log_file)
                    else:
                        print(f">> [PPT-PIPELINE] Agent 3 processing all Agent 1 data for fixes: {agent1_output}")
                        fixed_data = agent3.process_data(agent1_output, agent3_output, None)
                    
                    print(f">> [PPT-PIPELINE] Agent 3 completed: {agent3_output}")
                    
                    # Insert Agent 3's fixed data
                    if os.path.exists(agent3_output):
                        print(f">> [PPT-PIPELINE] Inserting Agent 3 fixed data into PostgreSQL...")
                        try:
                            # Check if Agent 3 has any data to insert
                            with open(agent3_output, 'r', encoding='utf-8') as f:
                                agent3_data = json.load(f)
                            
                            if agent3_data and len(agent3_data) > 0:
                                agent3_insertion_result = inserter.insert_json_file(agent3_output)
                                agent3_inserted = agent3_insertion_result.get('inserted_records', 0)
                                agent3_errors = agent3_insertion_result.get('errors', [])
                                
                                print(f">> [PPT-PIPELINE] Agent 3 insertion results:")
                                if agent3_inserted > 0:
                                    print(f"   - {agent3_inserted} records inserted successfully")
                                if agent3_errors:
                                    print(f"   - {len(agent3_errors)} insertion errors remain")
                                else:
                                    print(f"   - All Agent 3 records inserted successfully")
                            else:
                                print(f">> [PPT-PIPELINE] Agent 3 produced no records to insert")
                                
                        except Exception as e:
                            print(f">> [PPT-PIPELINE] Agent 3 database insertion failed: {e}")
                    
                except Exception as e:
                    print(f">> [PPT-PIPELINE] Agent 3 failed: {e}")
                    return {
                        'status': 'partial_success',
                        'error': f"Agent 3 failed: {e}",
                        'agent_failed': 3,
                        'processing_completed': True,
                        'agent1_output': agent1_output,
                        'insertion_result': insertion_result,
                        'insertion_errors_count': len(insertion_errors)
                    }
            else:
                print(f">> [PPT-PIPELINE] No errors found - skipping Agent 3")
                agent3_output = None
            
            # Final results
            print(f">> [PPT-PIPELINE] Modified pipeline completed successfully!")
            
            final_summary = {
                'status': 'success',
                'processing_completed': True,
                'pipeline_type': 'modified_ppt_pipeline_skip_agent2',
                'agent1_output': agent1_output,
                'agent3_output': agent3_output,
                'insertion_result': insertion_result,
                'insertion_errors_count': len(insertion_errors),
                'error_log_file': error_log_file,
                'total_records_inserted': (insertion_result.get('inserted_records', 0) if insertion_result else 0),
                'markdown_file': markdown_path,
                'source_powerpoint': source_file
            }
            
            # Add Agent 3 results if it was run
            if agent3_output and os.path.exists(agent3_output):
                final_summary['agent3_records_fixed'] = 'completed'
            
            return final_summary
                
        except Exception as e:
            print(f">> [PPT-PIPELINE] Exception in modified pipeline: {e}")
            import traceback
            print(f">> [PPT-PIPELINE] Full traceback:")
            traceback.print_exc()
            return {
                'status': 'exception',
                'error': str(e),
                'processing_completed': False,
                'markdown_file': markdown_path
            }
    
    def _combine_results(self, markdown_content: str, processed_tables: Dict, source_file: str) -> Dict[str, Any]:
        """Combine markdown content and processed tables into final result"""
        try:
            current_time = datetime.now().isoformat()
            
            # Process markdown into structured data
            markdown_data = []
            if markdown_content:
                # Split markdown by pages/sections
                sections = markdown_content.split('\n\n')
                current_section = ""
                
                for section in sections:
                    section = section.strip()
                    if section:
                        current_section += section + "\n\n"
                        
                        # If section is substantial, create an entry
                        if len(current_section.strip()) > 100:
                            markdown_data.append({
                                "source": f"{source_file}:markdown_content",
                                "title": self._extract_title_from_section(current_section),
                                "content": current_section.strip(),
                                "last_updated_date": current_time
                            })
                            current_section = ""
                
                # Add remaining content
                if current_section.strip():
                    markdown_data.append({
                        "source": f"{source_file}:markdown_content",
                        "title": "Additional Content",
                        "content": current_section.strip(),
                        "last_updated_date": current_time
                    })
            
            # Process table data
            table_data = []
            csv_files = []
            
            for slide_idx, slide_content in processed_tables.items():
                # Add table CSV entries
                for table in slide_content.get('tables', []):
                    csv_files.append({
                        "source": f"{source_file}:slide-{slide_idx}",
                        "title": f"Table from Slide {slide_idx}",
                        "csv_path": table.get('csv_path', ''),
                        "is_csv": True,
                        "last_updated_date": current_time
                    })
                
                # Add additional text from table slides to markdown
                additional_text = " ".join(slide_content.get('additional_text', []))
                if additional_text.strip():
                    markdown_data.append({
                        "source": f"{source_file}:slide-{slide_idx}",
                        "title": f"Text Content from Slide {slide_idx}",
                        "content": additional_text.strip(),
                        "last_updated_date": current_time
                    })
            
            # Combine all data
            all_data = markdown_data + csv_files
            
            return {
                "file_type": "powerpoint_text_extraction",
                "source_file": source_file,
                "processing_method": "text_extraction_with_gemini_vision",
                "total_entries": len(all_data),
                "markdown_entries": len(markdown_data),
                "csv_entries": len(csv_files),
                "last_updated_date": current_time,
                "data": all_data
            }
            
        except Exception as e:
            logger.error(f"Error combining results: {e}")
            return {
                "file_type": "powerpoint_text_extraction",
                "source_file": source_file,
                "error": str(e),
                "last_updated_date": datetime.now().isoformat(),
                "data": []
            }
    
    def _extract_title_from_section(self, section: str) -> str:
        """Extract title from a markdown section"""
        try:
            lines = section.split('\n')
            for line in lines:
                line = line.strip()
                if line and not line.startswith('!') and len(line) < 100:
                    # Clean the line to be a suitable title
                    title = re.sub(r'[#*\-]+', '', line).strip()
                    if title:
                        return title[:80]  # Limit title length
            
            # Fallback: use first 50 characters
            clean_text = re.sub(r'[#*\-!]+', '', section).strip()
            return clean_text[:50] + "..." if len(clean_text) > 50 else clean_text
            
        except:
            return "Content Section"


# Public interface function for pipeline integration
def process_powerpoint_with_text_extraction(file_path: Union[str, Path]) -> Dict[str, Any]:
    """
    Main entry point for PowerPoint processing with text extraction
    
    Args:
        file_path: Path to PowerPoint file
        
    Returns:
        Dictionary with processed content using text extraction workflow
    """
    processor = PowerPointProcessor()
    return processor.process_powerpoint_file(file_path)
